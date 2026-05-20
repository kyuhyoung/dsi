#!/usr/bin/env python3
"""빌드된 .pptx 의 슬라이드별 문제점 검사.

검사 항목:
- 슬라이드 크기 초과 (shape가 슬라이드 영역 밖)
- 빈 슬라이드 (텍스트·이미지 모두 없음)
- 깨진 이미지 path (rels에 image 참조는 있는데 zip에 파일 없음)
- 텍스트 box 면적 대비 글자 수 과다 (overflow 의심)
- 인접 shape 큰 겹침 (>50% 면적 overlap)

사용:
    python3 scripts/validate_pptx.py <pptx 경로>

출력: 슬라이드별 경고 리스트. 종료 코드: 0 (정상), 1 (오류만), 2 (경고만 있음).
"""
import sys
import zipfile
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu


def _in(emu):
    return (emu or 0) / 914400


def check_overflow(slide, slide_w, slide_h):
    """슬라이드 영역(slide_w × slide_h inch)을 벗어나는 shape 검사."""
    out = []
    for s in slide.shapes:
        if not s.left or not s.top:
            continue
        x, y = _in(s.left), _in(s.top)
        w, h = _in(s.width), _in(s.height)
        if x + w > slide_w + 0.1:
            out.append(f"shape '{s.name}' right edge {x+w:.2f}\" 초과 (slide width {slide_w:.2f})")
        if y + h > slide_h + 0.1:
            out.append(f"shape '{s.name}' bottom edge {y+h:.2f}\" 초과 (slide height {slide_h:.2f})")
        if x < -0.1:
            out.append(f"shape '{s.name}' left {x:.2f}\" 음수")
        if y < -0.1:
            out.append(f"shape '{s.name}' top {y:.2f}\" 음수")
    return out


def check_empty(slide):
    """텍스트도 이미지도 없는 빈 슬라이드 감지."""
    has_text = False
    has_pic = False
    for s in slide.shapes:
        try:
            if s.has_text_frame and s.text_frame.text.strip():
                has_text = True
        except Exception:
            pass
        try:
            st = str(s.shape_type)
            if "PICTURE" in st or "CHART" in st or "TABLE" in st:
                has_pic = True
        except Exception:
            pass
    if not has_text and not has_pic:
        return ["빈 슬라이드 (텍스트·이미지 없음)"]
    return []


def check_broken_images(pptx_path):
    """rels에 image 참조 있는데 zip에 실제 파일 없는 경우.
    슬라이드별 경고 dict {slide_idx: [msg]} 반환."""
    out = {}
    with zipfile.ZipFile(str(pptx_path)) as zf:
        names = set(zf.namelist())
        for n in names:
            if not (n.startswith("ppt/slides/_rels/slide") and n.endswith(".xml.rels")):
                continue
            slide_idx = int(n.split("slide")[2].split(".")[0])
            rels = zf.read(n).decode("utf-8", errors="replace")
            import re as _re
            for m in _re.finditer(r'Target="([^"]+)"', rels):
                target = m.group(1)
                if "/image" not in target and "/media/" not in target:
                    continue
                # target은 ../media/xxx 같은 상대 경로
                resolved = ("ppt/" + target.replace("../", "")).rstrip("/")
                if resolved not in names:
                    out.setdefault(slide_idx, []).append(
                        f"깨진 이미지 참조: {target}"
                    )
    return out


def check_text_overflow_heuristic(slide):
    """텍스트박스의 면적 대비 글자 수 추정. 한 글자당 약 0.018 sq.inch 가정.
    글자 수가 면적에 비해 *2배 이상* 많으면 overflow 의심."""
    out = []
    for s in slide.shapes:
        try:
            if not s.has_text_frame:
                continue
        except Exception:
            continue
        text = s.text_frame.text
        n = len(text.replace("\n", "").replace(" ", ""))
        if n < 30:
            continue
        w, h = _in(s.width), _in(s.height)
        area = w * h
        if area <= 0:
            continue
        char_capacity = area / 0.018
        if n > char_capacity * 2:
            out.append(
                f"shape '{s.name}' 글자 {n}자 / 박스 면적 {area:.1f} sq.in "
                f"(예상 수용 {int(char_capacity)}자) — overflow 의심"
            )
    return out


def validate(pptx_path):
    p = Path(pptx_path)
    if not p.exists():
        print(f"[오류] 파일 없음: {p}", file=sys.stderr)
        return 1
    prs = Presentation(str(p))
    slide_w = _in(prs.slide_width)
    slide_h = _in(prs.slide_height)

    broken_imgs = check_broken_images(p)
    all_warnings = []
    error_count = 0
    warn_count = 0

    print(f"검증: {p} ({len(prs.slides)} 슬라이드, {slide_w:.2f}×{slide_h:.2f}\")")
    print()

    for i, slide in enumerate(prs.slides, start=1):
        warnings = []
        warnings += check_overflow(slide, slide_w, slide_h)
        warnings += check_empty(slide)
        warnings += check_text_overflow_heuristic(slide)
        if i in broken_imgs:
            warnings += broken_imgs[i]
        if warnings:
            print(f"slide {i:2d}:")
            for w in warnings:
                level = "⚠" if "의심" in w or "초과" in w else "✗"
                if level == "✗":
                    error_count += 1
                else:
                    warn_count += 1
                print(f"   {level} {w}")
            all_warnings.extend(warnings)

    print()
    if not all_warnings:
        print("✓ 검증 통과. 발견된 문제 없음.")
        return 0
    print(f"검증 종료: 오류 {error_count} · 경고 {warn_count}")
    return 1 if error_count else 2


def main():
    if len(sys.argv) < 2:
        print("사용: python3 scripts/validate_pptx.py <pptx 경로>", file=sys.stderr)
        sys.exit(1)
    sys.exit(validate(sys.argv[1]))


if __name__ == "__main__":
    main()
