#!/usr/bin/env python3
"""PPT 구성안 YAML → pptx 변환.

v4: 견본 복제 + 시각화 함수 통합.
- 견본 슬라이드 복제 (다비오 템플릿)
- 시각요소(시각요소 항목)를 구조화된 type으로 자동 그리기:
  flow_arrow / matrix_2x2 / bar_chart / timeline / callout_cards
- 시각화 type이 있으면 본문 영역을 핵심메시지 상단 + 시각화 하단으로 분할

견본 매핑:
  종류=표지     → 견본 Slide 1
  종류=목차     → 견본 Slide 2
  종류=간지     → 견본 Slide 3
  종류=본문     → 견본 Slide 13
  종류=결론     → 견본 Slide 13
  종류=부록     → 견본 Slide 13
  종류=회사소개 → 견본 Slide 7
  종류=감사     → 견본 Slide 10
"""
import sys
import copy
import os
import zipfile
from pathlib import Path
import yaml
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR, MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE


STYLE_PATH = "templates/dabeeo_style.yaml"

# ── 기본값 (style.yaml 미존재 시 fallback) ──
TEMPLATE_PATH = "templates/Dabeeo_presentation_Template_2022_v1.pptx"

SAMPLE_MAP_1BASED = {
    "표지":     1,
    "목차":     2,
    "간지":     3,
    "회사소개": 7,
    "본문":     13,
    "결론":     13,
    "부록":     13,
    "감사":     10,
}

NAVY = RGBColor(0x1F, 0x36, 0x5C)
LIGHT_NAVY = RGBColor(0x4A, 0x6C, 0x9D)
GOLD = RGBColor(0xC8, 0x9B, 0x3C)
LIGHT_GOLD = RGBColor(0xE6, 0xC4, 0x7A)
GREY = RGBColor(0x70, 0x70, 0x70)
LIGHT_GREY = RGBColor(0xE8, 0xE8, 0xE8)
DARK = RGBColor(0x22, 0x22, 0x22)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# 표지 제목 자동 fit 룰 (style.yaml에서 override 가능)
TITLE_FIT_RULES = [
    {"max_chars": 10, "pt": 54},
    {"max_chars": 18, "pt": 42},
    {"max_chars": 28, "pt": 32},
    {"max_chars": 40, "pt": 26},
    {"max_chars": 999, "pt": 22},
]

# 본문 영역 — system_defaults.yaml 의 body_area 에서 로드 (style.yaml로 override 가능)
BODY_AREA = {}
COVER_TITLE_BOX = {"top": 3.0, "height": 2.7}
COVER_SUB1_BOX = {}  # left/top/width/height 옵션. 비어 있으면 본문 영역 width 사용
COVER_SUB2_BOX = {}
COVER_AUTO_FIT = False  # True 시 박스 dimension 기반 자동 폰트 계산
GRADIENT_WHITE_MASKS = []  # analyze_template이 감지한 흰 mask 좌표 (전 종류 견본, style.yaml 에서 로드)

# 시스템 default — templates/system_defaults.yaml 에서 로드 (필수).
# PY 안에 hardcoded default 없음. 값은 *전부 yaml*.
DEFAULTS = {}
EMPTY_BOX_MIN = {}
TYPOGRAPHY_FACTORS = {}
TYPOGRAPHY = {}
SHAPE_FINDING = {}
SHAPE_SAFETY = {}
SHAPE_DETECTION = {}
COVER_TYPOGRAPHY = {}
SECTION_DIVIDER_CFG = {}
THANKS_CFG = {}
TOC_CFG = {}
TOP_RIGHT_LABEL_PADDING = 0.05
TEXTBOX_MARGIN = {}
FIT_FONT_BOX = {}
CONTENT_LAYOUT_CFG = {}
BODY_TEXT_FONT_BY_LINES = []
ANALYZER_CFG = {}
KIND_REMAPPING = []
PLACEHOLDER_FIT_MAX_PT = {}
CHAPTER_MAPPING_POLICY = {}  # style.yaml의 chapter_mapping_policy
VISUAL_AREA_CLEANUP = {}  # system_defaults.yaml의 visual_area_cleanup
BAR_CHART_CFG = {}  # layout_vocabulary.yaml의 bar_chart 정책 (강조 색 등)

_SYSTEM_DEFAULTS_PATH = Path(__file__).parent.parent / "templates" / "system_defaults.yaml"

def _load_system_defaults():
    """시스템 전체 기본값 로드. 파일 없으면 시스템 죽음 (의도적)."""
    global DEFAULTS, EMPTY_BOX_MIN, TYPOGRAPHY_FACTORS, TYPOGRAPHY, FONT_BODY
    global SHAPE_FINDING, SHAPE_SAFETY, SHAPE_DETECTION, COVER_TYPOGRAPHY
    global SECTION_DIVIDER_CFG, THANKS_CFG, TOC_CFG, TOP_RIGHT_LABEL_PADDING
    global TEXTBOX_MARGIN, FIT_FONT_BOX, CONTENT_LAYOUT_CFG, BODY_TEXT_FONT_BY_LINES, ANALYZER_CFG
    global KIND_REMAPPING
    with open(_SYSTEM_DEFAULTS_PATH, "r", encoding="utf-8") as f:
        sd = yaml.safe_load(f) or {}
    DEFAULTS.update(sd.get("defaults", {}))
    EMPTY_BOX_MIN.update(sd.get("empty_box_min", {}))
    TYPOGRAPHY_FACTORS.update(sd.get("typography_factors", {}))
    TYPOGRAPHY.update(sd.get("typography", {}))
    if "font_body" in sd:
        FONT_BODY = sd["font_body"]
    SHAPE_FINDING.update(sd.get("shape_finding", {}))
    SHAPE_SAFETY.update(sd.get("shape_safety", {}))
    SHAPE_DETECTION.update(sd.get("shape_detection", {}))
    COVER_TYPOGRAPHY.update(sd.get("cover_typography", {}))
    SECTION_DIVIDER_CFG.update(sd.get("section_divider", {}))
    THANKS_CFG.update(sd.get("thanks", {}))
    TOC_CFG.update(sd.get("toc", {}))
    TOP_RIGHT_LABEL_PADDING = sd.get("top_right_label_padding", 0.05)
    TEXTBOX_MARGIN.update(sd.get("textbox_margin", {}))
    FIT_FONT_BOX.update(sd.get("fit_font_box", {}))
    CONTENT_LAYOUT_CFG.update(sd.get("content_layout", {}))
    BODY_TEXT_FONT_BY_LINES.extend(sd.get("body_text_font_by_lines", []))
    ANALYZER_CFG.update(sd.get("analyzer", {}))
    KIND_REMAPPING.clear()
    KIND_REMAPPING.extend(sd.get("kind_remapping", []) or [])
    PLACEHOLDER_FIT_MAX_PT.update(sd.get("placeholder_fit_max_pt", {}))
    VISUAL_AREA_CLEANUP.update(sd.get("visual_area_cleanup", {}))
    # 좌표 기본값 — style.yaml로 override 가능
    BODY_AREA.update(sd.get("body_area", {}))
    BODY_TEXTBOX_ANCHOR.update(sd.get("body_textbox_anchor", {}))
    cia = sd.get("company_intro_anchors", {}) or {}
    COMPANY_INTRO_ANCHORS["headline"].update(cia.get("headline", {}))
    COMPANY_INTRO_ANCHORS["sub"].update(cia.get("sub", {}))
CONTENT_LAYOUT = {"messages_height_with_visuals": 1.4, "visuals_top_gap": 0.2}
# system_defaults.yaml의 body_textbox_anchor, company_intro_anchors 에서 로드.
BODY_TEXTBOX_ANCHOR = {}
COMPANY_INTRO_ANCHORS = {"headline": {}, "sub": {}}
# 회사소개 슬라이드에 쓸 fill 함수 이름 (style.yaml company_intro.fill_function).
# 기본 'fill_company_intro' (전용 layout 시). 'fill_content'면 본문 layout 공유 시 사용.
COMPANY_INTRO_FILL_FN = "fill_company_intro"
PLACEHOLDER_IDX = {
    "표지": {"title": 0, "sub1": 1, "sub2": 2},
    "목차": {"title": 0, "nums_left": 1, "nums_right": 2, "titles_left": 3, "titles_right": 4},
    "본문": {"title": 0, "chap_no": 1, "chap_name": 2},
    "간지": {"title": 0, "chap_no": 1},
    "감사": {"title": 0, "footer": 11},
}

# 챕터 라벨 list — style.yaml의 chapter_labels 에서 로드.
# 라벨 종류·개수·위치·포맷·정렬 모두 yaml에 정의. PY는 단순 iterate.
# 빈 list면 라벨 안 그림 (layout 자리 없는 template용).
CHAPTER_LABELS = []

# 슬롯 매핑 룰 (style.yaml의 slot_finders)
# 비어 있으면 PLACEHOLDER_IDX fallback.
# 각 슬롯은 매핑 룰 리스트(우선순위 순으로 시도).
SLOT_FINDERS = {}  # type: dict[str, dict[str, list[dict]]]
FONT_BODY = ""  # 시스템 default yaml에서 로드. 회사 yaml override 가능.

# 시스템 default를 즉시 로드 (실패 시 시스템 죽음)
_load_system_defaults()

# 레이아웃 vocabulary 비율 (templates/layout_vocabulary.yaml에서 로드)
_LAYOUT_VOCAB_PATH = Path(__file__).parent.parent / "templates" / "layout_vocabulary.yaml"
def _load_layout_vocab():
    try:
        with open(_LAYOUT_VOCAB_PATH, "r", encoding="utf-8") as f:
            return yaml.safe_load(f) or {}
    except FileNotFoundError:
        return {}
_VOCAB = _load_layout_vocab()
LAYOUT_PARAMS = _VOCAB.get("layouts", {})
IMAGE_GRID_PATTERNS = {int(k): tuple(v) for k, v in (_VOCAB.get("image_grid_patterns", {})).items()}
# 시각요소 (visual) 렌더 파라미터
VISUALS_CFG = _VOCAB.get("visuals", {})
FLOW_ARROW_CFG = _VOCAB.get("flow_arrow", {})
MATRIX_2X2_CFG = _VOCAB.get("matrix_2x2", {})
TIMELINE_CFG = _VOCAB.get("timeline", {})
IMAGE_GRID_CFG = _VOCAB.get("image_grid", {})
CALLOUT_CARDS_CFG = _VOCAB.get("callout_cards", {})
# 변종 선택 정책 (시각요소 → 본문 견본 변종 매핑 룰 리스트)
CONTENT_VARIANT_PRIORITY = _VOCAB.get("content_variant_priority", [])
# bar_chart 정책 (강조 색 등)
BAR_CHART_CFG.update(_VOCAB.get("bar_chart", {}))


def _color_by_name(name, default=None):
    """yaml의 색 이름 ('gold', 'navy' 등) → RGBColor 객체.
    여러 곳에서 반복 박혔던 매핑을 단일 헬퍼로."""
    color_map = {
        "navy": NAVY, "light_navy": LIGHT_NAVY,
        "gold": GOLD, "light_gold": LIGHT_GOLD,
        "grey": GREY, "light_grey": LIGHT_GREY,
        "dark": DARK, "white": WHITE,
    }
    return color_map.get(name, default if default is not None else NAVY)

# 견본 복제 후처리 옵션 (style.yaml: sample_clone)
SAMPLE_CLONE = {
    # 견본의 placeholder가 아닌 텍스트박스 콘텐츠 자동 제거.
    # 회사 콘텐츠 박힌 templates용 (예: 새 다비오 V1).
    "clear_non_placeholder_text": False,
    # 보존할 키워드 (카피라이트, "Strictly Confidential" 등)
    "preserve_keywords": ["©", "Strictly", "STRICTLY", "Confidential", "CONFIDENTIAL", "보안 문서"],
}


def _hex_to_rgb(h):
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def load_style(style_path=STYLE_PATH):
    """style.yaml 로드하여 모듈 레벨 상수 갱신.
    파일 없거나 키 누락 시 기본값 유지 (fallback)."""
    global TEMPLATE_PATH, SAMPLE_MAP_1BASED, COMPANY_INTRO_FILL_FN
    global NAVY, LIGHT_NAVY, GOLD, LIGHT_GOLD, GREY, LIGHT_GREY, DARK, WHITE
    global TITLE_FIT_RULES, BODY_AREA, COVER_TITLE_BOX, CONTENT_LAYOUT, FONT_BODY

    p = Path(style_path)
    if not p.exists():
        print(f"style.yaml 없음 (기본값 사용): {style_path}", file=sys.stderr)
        return

    data = yaml.safe_load(p.read_text(encoding="utf-8")) or {}

    template = data.get("template", {})
    if "path" in template:
        TEMPLATE_PATH = template["path"]
    layouts = template.get("layouts")
    if isinstance(layouts, dict):
        SAMPLE_MAP_1BASED = dict(layouts)
    body = template.get("body_area")
    if isinstance(body, dict):
        BODY_AREA.update(body)

    colors = data.get("colors", {}) or {}
    if "navy" in colors: NAVY = _hex_to_rgb(colors["navy"])
    if "light_navy" in colors: LIGHT_NAVY = _hex_to_rgb(colors["light_navy"])
    if "gold" in colors: GOLD = _hex_to_rgb(colors["gold"])
    if "light_gold" in colors: LIGHT_GOLD = _hex_to_rgb(colors["light_gold"])
    if "grey" in colors: GREY = _hex_to_rgb(colors["grey"])
    if "light_grey" in colors: LIGHT_GREY = _hex_to_rgb(colors["light_grey"])
    if "dark" in colors: DARK = _hex_to_rgb(colors["dark"])
    if "white" in colors: WHITE = _hex_to_rgb(colors["white"])

    fonts = data.get("fonts", {}) or {}
    if "body" in fonts:
        FONT_BODY = fonts["body"]

    cover = data.get("cover", {}) or {}
    if isinstance(cover.get("title_fit"), list):
        TITLE_FIT_RULES = cover["title_fit"]
    if isinstance(cover.get("title_box"), dict):
        COVER_TITLE_BOX.update(cover["title_box"])
    if isinstance(cover.get("sub1_box"), dict):
        COVER_SUB1_BOX.update(cover["sub1_box"])
    if isinstance(cover.get("sub2_box"), dict):
        COVER_SUB2_BOX.update(cover["sub2_box"])
    global COVER_AUTO_FIT, GRADIENT_WHITE_MASKS
    if "auto_fit" in cover:
        COVER_AUTO_FIT = bool(cover["auto_fit"])
    # 위치 이동: cover.gradient_white_masks → sample_clone.gradient_white_masks
    sc_masks = (data.get("sample_clone", {}) or {}).get("gradient_white_masks")
    if isinstance(sc_masks, list):
        GRADIENT_WHITE_MASKS = sc_masks
    elif isinstance(cover.get("gradient_white_masks"), list):
        # 구버전 호환 (cover.gradient_white_masks 위치)
        GRADIENT_WHITE_MASKS = cover["gradient_white_masks"]

    # defaults (시스템 라벨) override
    defaults_cfg = data.get("defaults", {}) or {}
    if isinstance(defaults_cfg, dict):
        DEFAULTS.update(defaults_cfg)

    # empty_box_min (clear_non_placeholder_text 임계)
    ebm = (data.get("sample_clone", {}) or {}).get("empty_box_min")
    if isinstance(ebm, dict):
        EMPTY_BOX_MIN.update(ebm)

    # typography 추정 계수 + visual element font sizes
    typo = data.get("typography", {}) or {}
    if "char_width_factor" in typo:
        TYPOGRAPHY_FACTORS["char_width"] = typo["char_width_factor"]
    if "line_height_factor" in typo:
        TYPOGRAPHY_FACTORS["line_height"] = typo["line_height_factor"]
    # 모든 *_pt 키 자동 갱신
    for k, v in typo.items():
        if k.endswith("_pt") and isinstance(v, (int, float)):
            TYPOGRAPHY[k] = v

    content_cfg = data.get("content", {}) or {}
    if isinstance(content_cfg, dict):
        # body_textbox_anchor는 별도 dict로
        if "body_textbox_anchor" in content_cfg:
            BODY_TEXTBOX_ANCHOR.update(content_cfg["body_textbox_anchor"])
        # 나머지 키만 CONTENT_LAYOUT에
        for k, v in content_cfg.items():
            if k != "body_textbox_anchor":
                CONTENT_LAYOUT[k] = v

    ci = data.get("company_intro", {}) or {}
    if isinstance(ci, dict):
        if "headline_anchor" in ci:
            COMPANY_INTRO_ANCHORS["headline"].update(ci["headline_anchor"])
        if "sub_anchor" in ci:
            COMPANY_INTRO_ANCHORS["sub"].update(ci["sub_anchor"])
        if "fill_function" in ci:
            COMPANY_INTRO_FILL_FN = ci["fill_function"]

    ph_idx = data.get("placeholder_idx", {}) or {}
    if isinstance(ph_idx, dict):
        for kind, mapping in ph_idx.items():
            if kind in PLACEHOLDER_IDX and isinstance(mapping, dict):
                PLACEHOLDER_IDX[kind].update(mapping)
            elif isinstance(mapping, dict):
                PLACEHOLDER_IDX[kind] = dict(mapping)

    sc = data.get("sample_clone", {}) or {}
    if isinstance(sc, dict):
        SAMPLE_CLONE.update(sc)

    cls = data.get("chapter_labels", []) or []
    if isinstance(cls, list):
        CHAPTER_LABELS.clear()
        for item in cls:
            if isinstance(item, dict):
                CHAPTER_LABELS.append(item)

    cmp_ = data.get("chapter_mapping_policy", {}) or {}
    if isinstance(cmp_, dict):
        CHAPTER_MAPPING_POLICY.clear()
        CHAPTER_MAPPING_POLICY.update(cmp_)

    # section_divider 설정은 system_defaults + style.yaml override
    sdv = data.get("section_divider", {}) or {}
    if isinstance(sdv, dict):
        SECTION_DIVIDER_CFG.update(sdv)

    # slot_finders — 텍스트/위치 기반 슬롯 매핑 (placeholder 없는 템플릿용)
    sf = data.get("slot_finders", {}) or {}
    if isinstance(sf, dict):
        global SLOT_FINDERS
        SLOT_FINDERS = {}
        for kind, slots in sf.items():
            if isinstance(slots, dict):
                SLOT_FINDERS[kind] = {}
                for slot, rules in slots.items():
                    if isinstance(rules, list):
                        SLOT_FINDERS[kind][slot] = rules

    print(f"style 로드됨: {style_path}", file=sys.stderr)


def _remove_decorative_pictures(slide):
    """style.yaml의 sample_clone.remove_decorative_pictures 에 명시된 좌표
    매칭 도형(pic + connector + line) 제거. 분석기가 결정 → 빌더는 실행만.
    매칭 허용 오차: 위치 ±0.2 inch, 크기 ±0.3 inch."""
    rules = SAMPLE_CLONE.get("remove_decorative_pictures") or []
    if not rules:
        return
    for shape in list(slide.shapes):
        # placeholder는 건드리지 않음
        if shape.is_placeholder:
            continue
        sx = (shape.left or 0) / 914400
        sy = (shape.top or 0) / 914400
        sw = (shape.width or 0) / 914400
        sh = (shape.height or 0) / 914400
        for r in rules:
            if (abs(sx - r.get("left", 0)) < 0.2
                and abs(sy - r.get("top", 0)) < 0.2
                and abs(sw - r.get("width", 0)) < 0.3
                and abs(sh - r.get("height", 0)) < 0.3):
                sp = shape._element
                sp.getparent().remove(sp)
                break


def _clear_group_text(group_shape, preserve):
    """그룹 안에 중첩된 텍스트박스의 더미 텍스트를 비운다 (도형 자체는 보존).
    그룹 안 그룹도 재귀 처리. placeholder는 건드리지 않는다 (placeholder는 main loop에서 일괄 처리)."""
    for child in list(group_shape.shapes):
        try:
            st = child.shape_type
        except (ValueError, NotImplementedError):
            st = None
        if st == MSO_SHAPE_TYPE.GROUP:
            _clear_group_text(child, preserve)
            continue
        if child.is_placeholder:
            continue
        if child.has_text_frame:
            txt = child.text_frame.text.strip()
            if not txt:
                continue
            if any(kw in txt for kw in preserve):
                continue
            set_run_text(child, "")


def clear_non_placeholder_text(slide):
    """견본 복제 후 placeholder 아닌 텍스트/빈 큰 도형을 정리.
    (mask solidify·decorative picture 제거는 main loop에서 *flag 무관 항상* 호출됨)
    - 텍스트박스의 더미 텍스트 비우기 (그룹 내부도 재귀)
    - 콘텐츠 없는 큰 빈 도형(이미지 placeholder 자리) 제거
    style.yaml 의 sample_clone.preserve_keywords 에 매칭되는 텍스트는 보존."""
    preserve = SAMPLE_CLONE.get("preserve_keywords", []) or []
    # 큰 빈 도형 제거 임계값 (style.yaml의 sample_clone.empty_box_min 에서 로드)
    EMPTY_MIN_W = EMPTY_BOX_MIN.get("width", 1.5)
    EMPTY_MIN_H = EMPTY_BOX_MIN.get("height", 1.5)

    for shape in list(slide.shapes):
        if shape.is_placeholder:
            continue
        # shape type 식별
        try:
            st = shape.shape_type
        except (ValueError, NotImplementedError):
            st = None
        # 그룹은 도형 자체는 보존하되 *내부 텍스트는 비움*
        if st == MSO_SHAPE_TYPE.GROUP:
            _clear_group_text(shape, preserve)
            continue
        # 실제 그림·차트·표·미디어는 디자인 자산이므로 절대 제거 금지
        if st in (
            MSO_SHAPE_TYPE.PICTURE,
            MSO_SHAPE_TYPE.LINKED_PICTURE,
            MSO_SHAPE_TYPE.CHART,
            MSO_SHAPE_TYPE.TABLE,
            MSO_SHAPE_TYPE.MEDIA,
            MSO_SHAPE_TYPE.LINKED_OLE_OBJECT,
            MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT,
        ):
            continue
        w = (shape.width / 914400) if shape.width else 0
        h = (shape.height / 914400) if shape.height else 0
        if shape.has_text_frame:
            txt = shape.text_frame.text.strip()
            if not txt:
                # 빈 텍스트박스: 큰 도형이면 제거 (이미지 placeholder 자리 등)
                if w > EMPTY_MIN_W and h > EMPTY_MIN_H:
                    sp = shape._element
                    sp.getparent().remove(sp)
                continue
            if any(kw in txt for kw in preserve):
                continue
            set_run_text(shape, "")
        else:
            # 텍스트 프레임 없는 일반 도형 — 빈 큰 박스만 제거 (placeholder 자리)
            if w > EMPTY_MIN_W and h > EMPTY_MIN_H:
                sp = shape._element
                sp.getparent().remove(sp)


# ─────────────────────────────────────────
# 슬라이드 복제·유틸
# ─────────────────────────────────────────

_REL_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
_RID_ATTRS = (f"{{{_REL_NS}}}embed", f"{{{_REL_NS}}}link", f"{{{_REL_NS}}}id")
_SKIP_RELTYPES = (
    "/slideLayout",   # dest 가 이미 같은 layout 가짐
    "/slideMaster",
    "/notesSlide",    # 노트는 별도 처리 (없으면 생성 안 함)
)


def _remap_rids_in_element(element, rid_map):
    """XML 트리 내 r:embed / r:link / r:id 속성을 dest rId로 치환."""
    for el in element.iter():
        for attr in _RID_ATTRS:
            old_rid = el.get(attr)
            if old_rid and old_rid in rid_map:
                el.set(attr, rid_map[old_rid])


def _is_chapter_keyed(d):
    """dict의 모든 key가 int면 chapter-keyed로 판정 (2-level dict).
    {1: {기본: N}, 2: {기본: N}} vs {기본: N, 이미지: M}.
    """
    if not isinstance(d, dict) or not d:
        return False
    return all(isinstance(k, int) for k in d.keys())


def _apply_chapter_mapping(chap, available_chapters):
    """style.yaml의 chapter_mapping_policy에 따라 콘텐츠 장번호 → template 견본 챕터 매핑.
    available_chapters: sorted list of template 견본 챕터 번호.
    PY에 정책 박지 않음 — yaml의 type 값에 따라 단순 분기."""
    if not available_chapters:
        return None
    policy = CHAPTER_MAPPING_POLICY or {}
    ptype = policy.get("type", "1to1")
    if ptype == "manual":
        manual = policy.get("manual", {})
        return manual.get(chap)
    if ptype == "cycle":
        T = len(available_chapters)
        return available_chapters[(chap - 1) % T]
    if ptype == "last_repeat":
        if chap in available_chapters:
            return chap
        return available_chapters[-1]
    # 1to1 (default): chap이 available에 있으면 사용, 없으면 None
    return chap if chap in available_chapters else None


def _resolve_chapter_variants(sd, value):
    """value가 chapter-keyed (2-level) dict이면 slide의 장번호로 매핑 후 lookup.
    chapter-keyed 아니면 그대로 반환.
    매핑 정책은 yaml의 chapter_mapping_policy. 매칭 실패 시 명확한 오류."""
    if not _is_chapter_keyed(value):
        return value
    chap = sd.get("장번호")
    try:
        chap = int(chap) if chap is not None and str(chap).strip() != "" else None
    except (TypeError, ValueError):
        chap = None
    if chap is None:
        raise ValueError(
            f"슬라이드에 '장번호' 없음 — chapter-aware template은 장번호 필요. "
            f"슬라이드 제목: {sd.get('제목', '')!r}. "
            f"ppt-designer가 슬라이드에 *장번호*를 명시해야 합니다."
        )
    available = sorted(value.keys())
    mapped = _apply_chapter_mapping(chap, available)
    if mapped is None or mapped not in value:
        raise ValueError(
            f"장번호 {chap} → 매핑 실패 (정책: {(CHAPTER_MAPPING_POLICY or {}).get('type', '1to1')}). "
            f"available 챕터: {available}. "
            f"슬라이드 제목: {sd.get('제목', '')!r}"
        )
    return value[mapped]


def _normalize_variants(value):
    """layouts 값을 통일된 dict {tag: idx_1based}로 정규화.
    - int   → {"기본": int}
    - list  → {"기본": list[0], "변종1": list[1], ...}
    - dict  → 그대로
    - None  → None
    """
    if value is None:
        return None
    if isinstance(value, int):
        return {"기본": value}
    if isinstance(value, list):
        if not value:
            return None
        d = {"기본": value[0]}
        for i, v in enumerate(value[1:], start=1):
            d[f"변종{i}"] = v
        return d
    if isinstance(value, dict):
        return dict(value)
    return None


def _content_variant_priority(sd):
    """슬라이드 콘텐츠 특성 → 적합한 변종 태그 우선순위 리스트.
    정책 정의는 *templates/layout_vocabulary.yaml* 의 content_variant_priority.
    PY는 단순히 룰을 평가만 한다."""
    visuals = sd.get("시각요소", []) or []
    visual_types = {v.get("type") for v in visuals if isinstance(v, dict)}
    content = sd.get("내용", {}) or {}
    has_left_right = isinstance(content, dict) and ("좌" in content or "우" in content)

    def _rule_matches(rule):
        # 룰의 조건 키별 평가. 정확히 *어떤 조건*이 yaml에서 정의되었는지로 분기.
        if rule.get("always"):
            return True
        vt = rule.get("if_visual_type")
        if vt is not None:
            return vt in visual_types
        cci = rule.get("if_callout_card_items")
        if cci is not None:
            cards = [v for v in visuals
                     if isinstance(v, dict) and v.get("type") == "callout_cards"]
            return bool(cards and len(cards[0].get("items", []) or []) == cci)
        if rule.get("if_content_has_left_right") is True:
            return has_left_right
        return False

    priority = []
    for rule in CONTENT_VARIANT_PRIORITY:
        if _rule_matches(rule):
            tag = rule.get("variant")
            if tag and tag not in priority:
                priority.append(tag)
    return priority or ["기본"]


def pick_variant(sd, variants_dict):
    """콘텐츠 특성에 맞는 변종 1-based 인덱스 선택. 없으면 None."""
    if not variants_dict:
        return None
    priority = _content_variant_priority(sd)
    for tag in priority:
        if tag in variants_dict and variants_dict[tag] is not None:
            return variants_dict[tag]
    # 어느 태그도 매칭 안 되면 첫 항목
    for v in variants_dict.values():
        if v is not None:
            return v
    return None


def duplicate_slide(prs, src_slide):
    """src_slide 완전 복제 — shapes + 모든 relationships (이미지·차트·하이퍼링크 포함).

    이전 구현은 shape XML만 deepcopy하고 rels는 복사하지 않아 <p:pic> 등에서
    참조하는 rId가 dest 슬라이드 rels에 없어 이미지가 렌더링 안 됨.
    이 구현은 src.part.rels를 dest.part로 복사하면서 rId를 재매핑하고,
    deepcopy된 shape XML 내 r:embed/r:link/r:id를 새 rId로 치환한다.
    """
    blank = src_slide.slide_layout
    dest = prs.slides.add_slide(blank)

    # dest의 기본 placeholder 제거
    for shp in list(dest.shapes):
        sp = shp._element
        sp.getparent().remove(sp)

    # rels 복사 + src→dest rId 매핑.
    # src의 rId 번호 순서대로 등록해서 dest의 rId 분배가 src 의미를 보존하도록.
    # (LibreOffice 등 일부 렌더러가 동일 의미라도 rId 숫자 순서로 다르게 처리)
    rid_map = {}
    src_rels_sorted = sorted(
        src_slide.part.rels.values(),
        key=lambda r: int("".join(c for c in r.rId if c.isdigit()) or "0"),
    )
    for rel in src_rels_sorted:
        if any(rel.reltype.endswith(sk) for sk in _SKIP_RELTYPES):
            continue
        try:
            if rel.is_external:
                new_rid = dest.part.relate_to(
                    rel.target_ref, rel.reltype, is_external=True
                )
            else:
                new_rid = dest.part.relate_to(rel.target_part, rel.reltype)
            rid_map[rel.rId] = new_rid
        except Exception as e:
            print(
                f"warning: rel 복사 실패 {rel.rId} ({rel.reltype}): {e}",
                file=sys.stderr,
            )

    # shapes 복사 + rId 재매핑
    for shp in src_slide.shapes:
        el = shp._element
        new_el = copy.deepcopy(el)
        _remap_rids_in_element(new_el, rid_map)
        dest.shapes._spTree.append(new_el)

    return dest


def set_run_text(shape, new_text, keep_font=True):
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    saved_font = None
    if tf.paragraphs and tf.paragraphs[0].runs:
        r0 = tf.paragraphs[0].runs[0]
        saved_font = {
            "name": r0.font.name,
            "size": r0.font.size,
            "bold": r0.font.bold,
            "italic": r0.font.italic,
        }
        try:
            saved_font["color"] = r0.font.color.rgb
        except Exception:
            saved_font["color"] = None
    tf.text = ""
    lines = str(new_text).split("\n") if new_text is not None else [""]
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        r = p.add_run()
        r.text = line
        if keep_font and saved_font:
            if saved_font["name"]:
                r.font.name = saved_font["name"]
            if saved_font["size"]:
                r.font.size = saved_font["size"]
            if saved_font["bold"] is not None:
                r.font.bold = saved_font["bold"]
            if saved_font["italic"] is not None:
                r.font.italic = saved_font["italic"]
            if saved_font.get("color"):
                try:
                    r.font.color.rgb = saved_font["color"]
                except Exception:
                    pass


def set_run_font(run, name=None, size=14, bold=False, color=None):
    # name 생략 시 style.yaml에서 로드된 FONT_BODY 사용 (시스템 default)
    run.font.name = name or FONT_BODY
    run.font.size = Pt(size)
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color


def find_shape_by_text(slide, search):
    for shape in slide.shapes:
        if shape.has_text_frame:
            txt = shape.text_frame.text.strip()
            if txt.startswith(search) or search in txt:
                return shape
    return None


def find_shape_by_pos(slide, left_inch, top_inch, tol=None):
    """좌표 기반 shape 찾기.
    tol 미지정 시 system_defaults.yaml shape_finding.default_tolerance 사용.
    1차: slide-level shape 검사.
    2차: slide_layout shape 검사 — 매칭 시 slide-level로 deep copy 후 반환
         (layout shape는 read-only라 텍스트 교체용으로 슬라이드에 복제).
    이로써 layout-only slot(chap_no 등)도 빌더가 자동 처리."""
    if tol is None:
        tol = SHAPE_FINDING.get("default_tolerance", 0.3)
    target_l = Inches(left_inch)
    target_t = Inches(top_inch)
    # 1차: slide-level
    for shape in slide.shapes:
        if shape.left is None or shape.top is None:
            continue
        dl = abs(shape.left - target_l) / 914400
        dt = abs(shape.top - target_t) / 914400
        if dl < tol and dt < tol:
            return shape
    # 2차: layout-level — 매칭되면 slide-level로 deep copy
    try:
        layout = slide.slide_layout
    except Exception:
        return None
    for ls in layout.shapes:
        if ls.left is None or ls.top is None:
            continue
        if ls.is_placeholder:
            continue  # placeholder는 slide.shapes 통해 이미 검사됨
        dl = abs(ls.left - target_l) / 914400
        dt = abs(ls.top - target_t) / 914400
        if dl < tol and dt < tol:
            new_el = copy.deepcopy(ls._element)
            slide.shapes._spTree.append(new_el)
            # 새로 추가된 shape 객체 반환
            return slide.shapes[-1]
    return None


def find_shape_by_placeholder_idx(slide, idx):
    for shape in slide.shapes:
        if shape.is_placeholder and shape.placeholder_format.idx == idx:
            return shape
    return None


def find_shape_by_name(slide, name_substring):
    for shape in slide.shapes:
        if name_substring in shape.name:
            return shape
    return None


def find_shape_by_finders(slide, finders):
    """매핑 룰 리스트를 차례로 시도해 첫 매칭 shape 반환.

    각 룰은 dict, 다음 키 중 하나 사용:
    - {placeholder_idx: N}
    - {text_contains: "..."}
    - {position: {left: X, top: Y, tol: 0.3}}
    - {shape_name: "..."}
    """
    if not finders:
        return None
    for rule in finders:
        if not isinstance(rule, dict):
            continue
        if "placeholder_idx" in rule:
            try:
                idx = int(rule["placeholder_idx"])
            except (TypeError, ValueError):
                continue
            shape = find_shape_by_placeholder_idx(slide, idx)
            if shape is not None:
                return shape
        elif "text_contains" in rule:
            keyword = str(rule["text_contains"])
            shape = find_shape_by_text(slide, keyword)
            if shape is not None:
                return shape
        elif "position" in rule:
            pos = rule["position"]
            if isinstance(pos, dict) and "left" in pos and "top" in pos:
                shape = find_shape_by_pos(
                    slide, float(pos["left"]), float(pos["top"]),
                    tol=float(pos.get("tol", 0.3))
                )
                if shape is not None:
                    return shape
        elif "shape_name" in rule:
            shape = find_shape_by_name(slide, str(rule["shape_name"]))
            if shape is not None:
                return shape
    return None


def shapes_same(a, b):
    """두 shape이 같은 XML element를 가리키는지. python-pptx slide.shapes는
    매번 새 wrapper 객체를 반환하므로 id() 비교는 부정확. XML element 동일성으로 판단."""
    if a is None or b is None:
        return False
    try:
        return a._element is b._element
    except Exception:
        return a is b


def find_shape_by_slot(slide, kind, slot):
    """슬라이드 종류·슬롯 이름으로 shape 찾기.

    우선순위:
    1. SLOT_FINDERS[kind][slot] 의 매핑 룰 (slot_finders)
    2. PLACEHOLDER_IDX[kind][slot] 의 placeholder 인덱스 (fallback)
    """
    finders = SLOT_FINDERS.get(kind, {}).get(slot)
    if finders:
        shape = find_shape_by_finders(slide, finders)
        if shape is not None:
            return shape
    idx = PLACEHOLDER_IDX.get(kind, {}).get(slot)
    if idx is not None:
        return find_shape_by_placeholder_idx(slide, idx)
    return None


# ─────────────────────────────────────────
# 시각화 함수들 (NEW)
# ─────────────────────────────────────────

def draw_flow_arrow(slide, items, left, top, width, height):
    """N개 단계를 가로 화살표 흐름으로 그림. 파라미터는 layout_vocabulary.yaml flow_arrow 섹션."""
    n = len(items)
    if n == 0:
        return
    gap = FLOW_ARROW_CFG.get("gap", 0.2)
    box_h_ratio = FLOW_ARROW_CFG.get("box_height_ratio", 0.6)
    box_h_max = FLOW_ARROW_CFG.get("box_height_max_inch", 1.3)
    arr_neck = FLOW_ARROW_CFG.get("arrow_neck_width", 0.02)
    arr_head_w = FLOW_ARROW_CFG.get("arrow_head_width", 0.12)
    arr_head_l = FLOW_ARROW_CFG.get("arrow_head_length", 0.24)

    total_gap = gap * (n - 1)
    box_w = (width - total_gap) / n
    box_h = min(height * box_h_ratio, box_h_max)
    box_top = top + (height - box_h) / 2

    for i, item in enumerate(items):
        x = left + i * (box_w + gap)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(box_top), Inches(box_w), Inches(box_h)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = NAVY
        shape.line.fill.background()
        tf = shape.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = str(item)
        set_run_font(r, size=TYPOGRAPHY["flow_label_pt"], bold=True, color=WHITE)

        if i < n - 1:
            arr_x = x + box_w + arr_neck
            arr_y = box_top + box_h / 2 - arr_head_w
            arr_w = gap - 2 * arr_neck
            arr = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                Inches(arr_x), Inches(arr_y),
                Inches(arr_w), Inches(arr_head_l)
            )
            arr.fill.solid()
            arr.fill.fore_color.rgb = GOLD
            arr.line.fill.background()


def draw_matrix_2x2(slide, x_axis, y_axis, quadrants, left, top, width, height):
    """2×2 매트릭스. 파라미터는 layout_vocabulary.yaml matrix_2x2 섹션."""
    label_pad_left = MATRIX_2X2_CFG.get("label_pad_left", 0.4)
    label_pad_bottom = MATRIX_2X2_CFG.get("label_pad_bottom", 0.35)
    mx_left = left + label_pad_left
    mx_top = top + 0.05
    mx_w = width - label_pad_left - 0.05
    mx_h = height - label_pad_bottom - 0.1

    cell_w = mx_w / 2
    cell_h = mx_h / 2

    if isinstance(quadrants, dict):
        cells = [
            quadrants.get("Q2", quadrants.get("좌상", "")),
            quadrants.get("Q1", quadrants.get("우상", "")),
            quadrants.get("Q3", quadrants.get("좌하", "")),
            quadrants.get("Q4", quadrants.get("우하", "")),
        ]
    else:
        cells = list(quadrants) + [""] * (4 - len(quadrants))

    colors = [LIGHT_NAVY, NAVY, GREY, GOLD]
    positions = [
        (mx_left, mx_top),
        (mx_left + cell_w, mx_top),
        (mx_left, mx_top + cell_h),
        (mx_left + cell_w, mx_top + cell_h),
    ]

    for (x, y), label, color in zip(positions, cells, colors):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(y), Inches(cell_w), Inches(cell_h)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = WHITE
        shape.line.width = Pt(MATRIX_2X2_CFG.get("border_width_pt", 2))
        tf = shape.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        # label이 dict ({label, body, ...}) 이면 label\nbody 합쳐 그림.
        # 단순 문자열이면 그대로. schema yaml의 value_types 정의 따름.
        if isinstance(label, dict):
            head = str(label.get("label", "")).strip()
            body = str(label.get("body", "")).strip()
            label_text = f"{head}\n{body}".strip() if body else head
        else:
            label_text = str(label) if label is not None else ""
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = label_text
        set_run_font(r, size=TYPOGRAPHY["matrix_axis_pt"], bold=True, color=WHITE)

    # x축 라벨 (하단 중앙)
    x_lbl = slide.shapes.add_textbox(
        Inches(mx_left), Inches(mx_top + mx_h + 0.05),
        Inches(mx_w), Inches(label_pad_bottom)
    )
    tf = x_lbl.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = f"→  {x_axis}"
    set_run_font(r, size=TYPOGRAPHY["matrix_title_pt"], bold=True, color=DARK)

    # y축 라벨 (좌측 상단)
    y_lbl_offset = MATRIX_2X2_CFG.get("y_axis_label_offset", 0.05)
    y_lbl_w = MATRIX_2X2_CFG.get("y_axis_label_width", 0.4)
    y_lbl = slide.shapes.add_textbox(
        Inches(left), Inches(mx_top),
        Inches(label_pad_left - y_lbl_offset), Inches(y_lbl_w)
    )
    tf = y_lbl.text_frame
    r = tf.paragraphs[0].add_run()
    r.text = f"↑\n{y_axis}"
    set_run_font(r, size=TYPOGRAPHY["matrix_text_pt"], bold=True, color=DARK)


def draw_bar_chart(slide, title, categories, series, left, top, width, height,
                   highlight_indices=None):
    """막대 차트. highlight_indices에 포함된 카테고리는 강조 색 (GOLD).
    yaml의 highlight_indices 또는 highlight_categories 로 지정. 빌더에서 list[int]로 정규화."""
    if not categories or not series:
        return
    highlight_set = set(highlight_indices or [])
    chart_data = CategoryChartData()
    chart_data.categories = list(categories)
    for name, values in series.items():
        chart_data.add_series(str(name), list(values))
    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(left), Inches(top),
        Inches(width), Inches(height),
        chart_data
    )
    chart = chart_shape.chart
    if title:
        chart.has_title = True
        chart.chart_title.text_frame.text = str(title)
        for run in chart.chart_title.text_frame.paragraphs[0].runs:
            set_run_font(run, size=TYPOGRAPHY["chart_label_pt"], bold=True, color=NAVY)
    chart.has_legend = len(series) > 1

    # 강조 처리 — 각 시리즈의 *highlight_indices* 위치 데이터 포인트만 GOLD.
    # python-pptx는 데이터 포인트 색 직접 설정 어려움 → XML 직접 조작.
    if highlight_set:
        for plot in chart.plots:
            for ser in plot.series:
                _highlight_series_points(ser, highlight_set)


def _highlight_series_points(series, highlight_indices):
    """막대 차트의 시리즈 안에서 *특정 인덱스 데이터 포인트*에 강조 색 적용.
    강조 색은 layout_vocabulary.yaml의 bar_chart.highlight_color (style.yaml colors 키).
    python-pptx 공식 API 부족으로 lxml XML 직접 조작 (알고리즘 — yaml로 표현 불가)."""
    from lxml import etree
    ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    ns_c = "http://schemas.openxmlformats.org/drawingml/2006/chart"
    # 강조 색 결정 — yaml 정책
    color = _color_by_name(BAR_CHART_CFG.get("highlight_color", "gold"), default=GOLD)
    ser_el = series._element
    for old in ser_el.findall(f"{{{ns_c}}}dPt"):
        ser_el.remove(old)
    for idx in highlight_indices:
        dPt = etree.SubElement(ser_el, f"{{{ns_c}}}dPt")
        idx_el = etree.SubElement(dPt, f"{{{ns_c}}}idx")
        idx_el.set("val", str(idx))
        invertIfNegative = etree.SubElement(dPt, f"{{{ns_c}}}invertIfNegative")
        invertIfNegative.set("val", "0")
        bubble3D = etree.SubElement(dPt, f"{{{ns_c}}}bubble3D")
        bubble3D.set("val", "0")
        spPr = etree.SubElement(dPt, f"{{{ns_c}}}spPr")
        solidFill = etree.SubElement(spPr, f"{{{ns_a}}}solidFill")
        srgbClr = etree.SubElement(solidFill, f"{{{ns_a}}}srgbClr")
        srgbClr.set("val", f"{color}")


def draw_timeline(slide, items, left, top, width, height):
    """가로 타임라인 (점·라인·라벨). 파라미터는 layout_vocabulary.yaml timeline 섹션."""
    n = len(items)
    if n == 0:
        return
    line_y = top + height * 0.5
    pad_x = TIMELINE_CFG.get("pad_x", 0.4)
    line_h = TIMELINE_CFG.get("line_height", 0.04)
    dot_size = TIMELINE_CFG.get("dot_size_inch", 0.36)

    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left + pad_x), Inches(line_y - line_h / 2),
        Inches(width - 2 * pad_x), Inches(line_h)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = NAVY
    line.line.fill.background()

    if n == 1:
        step = 0
        start_x = left + width / 2
    else:
        step = (width - 2 * pad_x) / (n - 1)
        start_x = left + pad_x

    for i, item in enumerate(items):
        x = start_x + i * step
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(x - dot_size / 2), Inches(line_y - dot_size / 2),
            Inches(dot_size), Inches(dot_size)
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = GOLD
        dot.line.color.rgb = NAVY
        dot.line.width = Pt(2)

        if isinstance(item, dict):
            label = item.get("label", "")
            content = item.get("content", "")
        else:
            label = str(item)
            content = ""

        # 라벨 배치 (홀수=위, 짝수=아래)
        if i % 2 == 0:
            lbl_top = top
            lbl_h = (line_y - top) - 0.25
        else:
            lbl_top = line_y + 0.25
            lbl_h = (top + height - line_y) - 0.25

        lbl_w = min(step * 0.95 if step else 2.0, 2.0)
        lbl_box = slide.shapes.add_textbox(
            Inches(x - lbl_w / 2), Inches(lbl_top),
            Inches(lbl_w), Inches(max(lbl_h, 0.5))
        )
        tf = lbl_box.text_frame
        tf.word_wrap = True
        if i % 2 == 0:
            tf.vertical_anchor = MSO_ANCHOR.BOTTOM
        else:
            tf.vertical_anchor = MSO_ANCHOR.TOP

        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = label
        set_run_font(r, size=TYPOGRAPHY["timeline_label_pt"], bold=True, color=NAVY)
        if content:
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.CENTER
            r2 = p2.add_run()
            r2.text = content
            set_run_font(r2, size=TYPOGRAPHY["timeline_content_pt"], color=GREY)


def draw_image(slide, path, left, top, width, height, fit="contain", caption=None):
    """이미지 삽입. caption 있으면 이미지 아래 작은 텍스트로 추가.
    fit=contain: 비율 유지하며 영역 안에 맞춤 (기본)
    fit=cover:   영역을 가득 채움 (잘릴 수 있음)
    fit=exact:   영역에 강제로 채움 (왜곡 가능)"""
    import os
    if not os.path.exists(path):
        box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_GREY
        box.line.color.rgb = GREY
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = f"[이미지 없음]\n{path}"
        set_run_font(r, size=TYPOGRAPHY["caption_pt"], color=GREY)
        return None

    # 캡션 영역 확보 (있을 때만)
    caption_h = 0.3 if caption else 0
    img_area_h = max(0.3, height - caption_h)

    try:
        from PIL import Image as PILImage
        img = PILImage.open(path)
        iw, ih = img.size
        aspect = iw / ih
    except Exception:
        aspect = 16 / 9

    target_aspect = width / img_area_h if img_area_h > 0 else 1.0

    if fit == "contain":
        if aspect > target_aspect:
            w = width
            h = width / aspect
            l = left
            t = top + (img_area_h - h) / 2
        else:
            h = img_area_h
            w = img_area_h * aspect
            l = left + (width - w) / 2
            t = top
    elif fit == "cover":
        if aspect > target_aspect:
            h = img_area_h
            w = img_area_h * aspect
            l = left - (w - width) / 2
            t = top
        else:
            w = width
            h = width / aspect
            l = left
            t = top - (h - img_area_h) / 2
    else:
        l, t, w, h = left, top, width, img_area_h

    slide.shapes.add_picture(
        path, Inches(l), Inches(t), Inches(w), Inches(h)
    )

    # 캡션 출력 (있을 때)
    if caption:
        cap_box = slide.shapes.add_textbox(
            Inches(left), Inches(top + img_area_h),
            Inches(width), Inches(caption_h)
        )
        cap_tf = cap_box.text_frame
        cap_tf.word_wrap = True
        cap_p = cap_tf.paragraphs[0]
        cap_p.alignment = PP_ALIGN.CENTER
        cap_r = cap_p.add_run()
        cap_r.text = str(caption)
        set_run_font(cap_r, size=TYPOGRAPHY["caption_pt"], color=GREY)
        cap_r.font.italic = True


def draw_mermaid(slide, mermaid_code, left, top, width, height, theme="default", background="white"):
    """Mermaid 텍스트를 PNG로 변환해서 슬라이드에 삽입.
    캐시: .cache/mermaid/<hash>.png로 같은 코드는 재생성 안 함."""
    import subprocess
    import hashlib

    if not mermaid_code or not mermaid_code.strip():
        return

    cache_key = hashlib.md5(
        (mermaid_code + theme + background).encode("utf-8")
    ).hexdigest()[:12]
    cache_dir = Path(".cache/mermaid")
    cache_dir.mkdir(parents=True, exist_ok=True)
    png_path = cache_dir / f"{cache_key}.png"

    if not png_path.exists():
        mmd_path = cache_dir / f"{cache_key}.mmd"
        mmd_path.write_text(mermaid_code, encoding="utf-8")

        mmdc_bin = Path("node_modules/.bin/mmdc")
        if not mmdc_bin.exists():
            # mmdc 없으면 placeholder
            box = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(left), Inches(top), Inches(width), Inches(height)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = LIGHT_GREY
            tf = box.text_frame
            tf.word_wrap = True
            r = tf.paragraphs[0].add_run()
            r.text = "[Mermaid CLI 미설치: node_modules/.bin/mmdc 없음]"
            set_run_font(r, size=11, color=GREY)
            return

        try:
            # 고해상도 PNG: -s 3 (스케일 3배), -w/-H로 베이스 픽셀 지정
            result = subprocess.run(
                [str(mmdc_bin), "-i", str(mmd_path), "-o", str(png_path),
                 "-t", theme, "-b", background,
                 "-w", "1600", "-H", "900", "-s", "2", "--quiet"],
                capture_output=True, text=True, timeout=90
            )
            if result.returncode != 0:
                print(f"mmdc 오류: {result.stderr[:300]}", file=sys.stderr)
                return
        except Exception as e:
            print(f"mmdc 실행 실패: {e}", file=sys.stderr)
            return

    if png_path.exists():
        draw_image(slide, str(png_path), left, top, width, height, fit="contain")


def draw_callout_cards(slide, items, left, top, width, height):
    """N개 강조 카드 가로 배치. 파라미터는 layout_vocabulary.yaml callout_cards 섹션."""
    n = len(items)
    if n == 0:
        return
    gap = IMAGE_GRID_CFG.get("gap", 0.25)  # 같은 grid 간격 사용
    bar_h_margin = CALLOUT_CARDS_CFG.get("h_margin", 0.12)
    bar_v_margin = CALLOUT_CARDS_CFG.get("v_margin", 0.20)
    bar_thickness = CALLOUT_CARDS_CFG.get("bar_thickness", 0.06)

    total_gap = gap * (n - 1)
    card_w = (width - total_gap) / n
    target_card_h = height
    card_top = top

    for i, item in enumerate(items):
        x = left + i * (card_w + gap)
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(card_top), Inches(card_w), Inches(target_card_h)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = NAVY
        card.line.width = Pt(1.5)

        # 상단 강조 라인 (카드 둥근 모서리 안쪽에 끼움)
        bar = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x + bar_v_margin),
            Inches(card_top + bar_h_margin),
            Inches(card_w - 2 * bar_v_margin),
            Inches(bar_thickness)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = GOLD
        bar.line.fill.background()

        if isinstance(item, dict):
            title = item.get("title", "")
            body = item.get("body", "")
        else:
            title = str(item)
            body = ""

        # 콘텐츠를 카드 정중앙에 명시 위치로 배치 (vertical_anchor는 LibreOffice 변환에서 불안정)
        estimated_content_h = CALLOUT_CARDS_CFG.get("estimated_content_height", 1.2)
        content_top = card_top + (target_card_h - estimated_content_h) / 2
        text_box = slide.shapes.add_textbox(
            Inches(x + 0.25), Inches(content_top),
            Inches(card_w - 0.5), Inches(estimated_content_h + 0.3)
        )
        tf = text_box.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = f"0{i + 1}"
        set_run_font(r, size=TYPOGRAPHY["callout_number_pt"], bold=True, color=GOLD)
        p.space_after = Pt(6)

        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = title
        set_run_font(r2, size=TYPOGRAPHY["callout_title_pt"], bold=True, color=NAVY)
        p2.space_after = Pt(6)

        if body:
            p3 = tf.add_paragraph()
            r3 = p3.add_run()
            r3.text = body
            set_run_font(r3, size=TYPOGRAPHY["callout_body_pt"], color=DARK)


def _resolve_highlight_indices(v):
    """yaml의 highlight_indices (list[int]) 또는 highlight_categories (list[str]) → list[int].
    둘 다 없으면 빈 list."""
    if "highlight_indices" in v:
        return [int(i) for i in v["highlight_indices"] if isinstance(i, int) or str(i).isdigit()]
    if "highlight_categories" in v:
        names = v["highlight_categories"]
        cats = v.get("categories", [])
        return [cats.index(n) for n in names if n in cats]
    return []


VISUAL_DISPATCH = {
    "flow_arrow": lambda s, v, l, t, w, h: draw_flow_arrow(s, v.get("items", []), l, t, w, h),
    "matrix_2x2": lambda s, v, l, t, w, h: draw_matrix_2x2(
        s, v.get("x_axis", ""), v.get("y_axis", ""), v.get("quadrants", {}), l, t, w, h),
    "bar_chart": lambda s, v, l, t, w, h: draw_bar_chart(
        s, v.get("title", ""), v.get("categories", []), v.get("series", {}), l, t, w, h,
        highlight_indices=_resolve_highlight_indices(v)),
    "timeline": lambda s, v, l, t, w, h: draw_timeline(s, v.get("items", []), l, t, w, h),
    "callout_cards": lambda s, v, l, t, w, h: draw_callout_cards(s, v.get("items", []), l, t, w, h),
    "image": lambda s, v, l, t, w, h: draw_image(
        s, v.get("path", ""), l, t, w, h,
        fit=v.get("fit", "contain"), caption=v.get("caption")),
    "mermaid": lambda s, v, l, t, w, h: draw_mermaid(
        s, v.get("code", ""), l, t, w, h,
        theme=v.get("theme", "default"),
        background=v.get("background", "white")),
}


def _remove_listed_shapes(slide):
    """style.yaml의 sample_clone.remove_shape_coords 좌표 매칭 shape 제거.
    *알고리즘만*: 좌표 비교 + shape 제거. 임계·키워드 0.
    """
    coords_list = SAMPLE_CLONE.get("remove_shape_coords", []) or []
    if not coords_list:
        return
    tol = SAMPLE_CLONE.get("shape_coords_tol", 0.1)
    for shape in list(slide.shapes):
        if shape.is_placeholder:
            continue
        if shape.left is None or shape.top is None:
            continue
        sl = shape.left / 914400
        st = shape.top / 914400
        sw = (shape.width or 0) / 914400
        sh = (shape.height or 0) / 914400
        for spec in coords_list:
            cl = spec.get("left", 0)
            ct = spec.get("top", 0)
            # 좌표 매칭 (width·height는 옵션 — 명시되어 있으면 추가 검사)
            if abs(sl - cl) > tol or abs(st - ct) > tol:
                continue
            if "width" in spec and abs(sw - spec["width"]) > tol:
                continue
            if "height" in spec and abs(sh - spec["height"]) > tol:
                continue
            sp = shape._element
            parent = sp.getparent()
            if parent is not None:
                parent.remove(sp)
            break


def _clear_shapes_in_area(slide, left, top, width, height, keep_ids=None):
    """주어진 직사각형 영역과 겹치는 *견본 잔재 도형* 제거.
    정책 (preserve_shape_types, picture_preserve_min_*) 은 system_defaults.yaml의 visual_area_cleanup.
    제거 비대상: keep_ids에 포함된 shape + yaml에 정의된 보존 type.
    """
    keep_ids = keep_ids or set()
    cleanup_cfg = VISUAL_AREA_CLEANUP
    preserve_type_names = set(cleanup_cfg.get("preserve_shape_types", []) or [])
    # MSO_SHAPE_TYPE enum 값 매핑
    preserve_enums = set()
    for name in preserve_type_names:
        try:
            preserve_enums.add(getattr(MSO_SHAPE_TYPE, name))
        except AttributeError:
            pass
    pic_min_w = cleanup_cfg.get("picture_preserve_min_width", 1.5)
    pic_min_h = cleanup_cfg.get("picture_preserve_min_height", 1.5)

    area_l, area_t = left, top
    area_r, area_b = left + width, top + height
    for shape in list(slide.shapes):
        if id(shape) in keep_ids:
            continue
        if shape.is_placeholder:
            continue
        try:
            st = shape.shape_type
        except (ValueError, NotImplementedError):
            st = None
        if st in preserve_enums:
            continue
        if shape.left is None or shape.top is None or shape.width is None or shape.height is None:
            continue
        sl = shape.left / 914400
        st_ = shape.top / 914400
        sr = sl + shape.width / 914400
        sb = st_ + shape.height / 914400
        if sr < area_l or sl > area_r or sb < area_t or st_ > area_b:
            continue
        w_in = shape.width / 914400
        h_in = shape.height / 914400
        # 큰 PICTURE는 콘텐츠로 간주해 보존 (yaml의 임계)
        if st in (MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.LINKED_PICTURE):
            if w_in > pic_min_w and h_in > pic_min_h:
                continue
        sp = shape._element
        parent = sp.getparent()
        if parent is not None:
            parent.remove(sp)


def render_visuals(slide, visuals, left, top, width, height):
    """구조화 시각요소 리스트 배치. 이미지가 여러 장이면 가로 그리드.
    이미지 + 비-이미지 혼재 시 비-이미지(상단) + 이미지 그리드(하단) 분할.
    시각요소 그리기 전에 *해당 영역의 견본 잔재 도형 제거*.
    """
    structured = [v for v in visuals if isinstance(v, dict) and "type" in v]
    if not structured:
        return False

    # 시각요소 영역의 견본 도형 정리 (line·dot·작은 picture 등 디자인 잔재)
    _clear_shapes_in_area(slide, left, top, width, height)

    images = [v for v in structured if v.get("type") == "image"]
    others = [v for v in structured if v.get("type") != "image"]
    gap = VISUALS_CFG.get("gap_between", 0.15)

    # case A: 비-이미지만
    if not images:
        return _render_vertical(slide, others, left, top, width, height, gap)

    # case B: 이미지만 (여러 장이면 가로 그리드)
    if not others:
        return _render_image_grid(slide, images, left, top, width, height, gap)

    # case C: 혼재 — 상단 비-이미지 + 하단 이미지. 비율은 layout_vocabulary.yaml visuals.mixed_upper_ratio.
    upper_h = height * VISUALS_CFG.get("mixed_upper_ratio", 0.4)
    lower_h = height - upper_h - gap
    _render_vertical(slide, others, left, top, width, upper_h, gap)
    _render_image_grid(slide, images, left, top + upper_h + gap,
                       width, lower_h, gap)
    return True


def _render_vertical(slide, items, left, top, width, height, gap):
    if not items:
        return False
    n = len(items)
    each_h = (height - gap * (n - 1)) / n
    for i, v in enumerate(items):
        v_top = top + i * (each_h + gap)
        fn = VISUAL_DISPATCH.get(v.get("type"))
        if fn:
            fn(slide, v, left, v_top, width, each_h)
    return True


def _render_image_grid(slide, images, left, top, width, height, gap):
    """이미지 N장 그리드 배치. (rows, cols)는 layout_vocabulary.yaml 의
    image_grid_patterns 에서 로드."""
    n = len(images)
    rows, cols = IMAGE_GRID_PATTERNS.get(n, IMAGE_GRID_PATTERNS.get(6, (2, 3)))

    cell_w = (width - gap * (cols - 1)) / cols
    cell_h = (height - gap * (rows - 1)) / rows
    for i, v in enumerate(images):
        if i >= rows * cols:
            break
        r = i // cols
        c = i % cols
        v_left = left + c * (cell_w + gap)
        v_top = top + r * (cell_h + gap)
        draw_image(
            slide, v.get("path", ""), v_left, v_top, cell_w, cell_h,
            fit=v.get("fit", "contain"), caption=v.get("caption"),
        )
    return True


def draw_messages_box(slide, messages, contents_dict, left, top, width, height):
    """핵심메시지·내용을 textbox로 그림 (시각요소가 있을 때 상단 영역)."""
    if not messages and not contents_dict:
        return None
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for m in messages:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        r = p.add_run()
        r.text = f"▶  {m}"
        set_run_font(r, size=TYPOGRAPHY["body_message_pt"], bold=True, color=NAVY)
        p.space_after = Pt(12)

    if contents_dict and isinstance(contents_dict, dict):
        for k, v in contents_dict.items():
            p = tf.add_paragraph() if not first else tf.paragraphs[0]
            first = False
            if isinstance(v, list):
                r = p.add_run()
                r.text = f"■ {k}:"
                set_run_font(r, size=TYPOGRAPHY["content_pt"], bold=True, color=DARK)
                p.space_after = Pt(6)
                for sub in v:
                    sp = tf.add_paragraph()
                    sr = sp.add_run()
                    sr.text = f"    · {sub}"
                    set_run_font(sr, size=TYPOGRAPHY["sub_item_pt"], color=GREY)
                    sp.space_after = Pt(4)
            else:
                r = p.add_run()
                r.text = f"■ {k}: {v}"
                set_run_font(r, size=TYPOGRAPHY["content_pt"], color=DARK)
                p.space_after = Pt(8)
    return box


# ─────────────────────────────────────────
# fill 함수들
# ─────────────────────────────────────────

def fit_title_size(text):
    """표지 제목 길이 기반 폰트 크기 (TITLE_FIT_RULES fallback)."""
    n = len(text.replace(' ', '').replace('\n', ''))
    for rule in TITLE_FIT_RULES:
        if n <= rule["max_chars"]:
            return rule["pt"]
    return TITLE_FIT_RULES[-1]["pt"]


def fit_font_to_box(text, box_w, box_h, max_pt=None, min_pt=None):
    """텍스트가 box 안에 word_wrap으로 들어가는 가장 큰 font pt 결정 (binary search).
    한글·영문 혼합 가정으로 character width 0.9 * pt/72, line height 1.25 * pt/72 추정.
    char 단위 줄바꿈으로 보수적 예상(실제 자동 줄바꿈은 단어 단위라 더 잘 들어감).
    기본 범위는 system_defaults.yaml fit_font_box.{max_pt, min_pt}.
    """
    if max_pt is None:
        max_pt = FIT_FONT_BOX.get("max_pt", 60)
    if min_pt is None:
        min_pt = FIT_FONT_BOX.get("min_pt", 12)
    if not text or box_w <= 0 or box_h <= 0:
        return max_pt
    explicit_lines = str(text).split("\n")
    lo, hi = min_pt, max_pt
    best = min_pt
    while lo <= hi:
        mid = (lo + hi) // 2
        char_w = mid / 72.0 * TYPOGRAPHY_FACTORS["char_width"]
        line_h = mid / 72.0 * TYPOGRAPHY_FACTORS["line_height"]
        chars_per_line = max(1, int(box_w / char_w))
        total_lines = sum(
            max(1, (len(line) + chars_per_line - 1) // chars_per_line)
            for line in explicit_lines
        )
        if total_lines * line_h <= box_h:
            best = mid
            lo = mid + 1
        else:
            hi = mid - 1
    return best


def _add_cover_textbox(slide, text, box_cfg, default_pt=24, color=None, bold=True, align=None):
    """fallback textbox 생성 (cover title/sub fallback용)."""
    if not text:
        return
    left = box_cfg.get("left", 1.05)
    top = box_cfg.get("top", 3.0)
    width = box_cfg.get("width", 11.20)
    height = box_cfg.get("height", 0.8)
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if align is not None:
        p.alignment = align
    r = p.add_run()
    r.text = str(text)
    set_run_font(r, size=default_pt, bold=bold, color=color or NAVY)


def _force_white_fill(shape):
    """텍스트박스의 배경을 명시 흰 solid fill로 강제.
    표지 견본의 그라데이션 mask를 LibreOffice가 비결정적으로 처리하는
    의존성 차단. 콘텐츠 텍스트가 그림 위에 비치는 현상 방지."""
    try:
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    except Exception as e:
        print(f"warning: fill 강제 실패 ({getattr(shape, 'name', '?')}): {e}", file=sys.stderr)


def _solidify_gradient_masks(slide):
    """style.yaml의 sample_clone.gradient_white_masks 좌표 매칭 도형의 fill을 *yaml이 명시한 색*으로 solid 변환.
    판단(어느 좌표·어느 색)은 analyze_template이 결정 → 빌더는 실행만.
    매칭 허용 오차: 위치 ±0.2 inch, 크기 ±0.3 inch.
    *전 종류 견본 공통* 처리, 어느 회사 템플릿이든 작동."""
    rules = GRADIENT_WHITE_MASKS or []
    if not rules:
        return
    A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
    P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
    from lxml import etree as _etree

    sp_elements = []
    def _collect(el):
        for child in el:
            tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
            if tag == "sp":
                sp_elements.append(child)
            elif tag == "grpSp":
                _collect(child)
    _collect(slide.shapes._spTree)

    for sp in sp_elements:
        sp_pr = sp.find(f"{{{P_NS}}}spPr")
        if sp_pr is None:
            continue
        xfrm = sp_pr.find(f"{{{A_NS}}}xfrm")
        if xfrm is None:
            continue
        off = xfrm.find(f"{{{A_NS}}}off")
        ext = xfrm.find(f"{{{A_NS}}}ext")
        if off is None or ext is None:
            continue
        sx = int(off.get("x", 0)) / 914400
        sy = int(off.get("y", 0)) / 914400
        sw = int(ext.get("cx", 0)) / 914400
        sh = int(ext.get("cy", 0)) / 914400
        matched_rule = None
        for r in rules:
            if (abs(sx - r.get("left", 0)) < 0.2
                and abs(sy - r.get("top", 0)) < 0.2
                and abs(sw - r.get("width", 0)) < 0.3
                and abs(sh - r.get("height", 0)) < 0.3):
                matched_rule = r
                break
        if matched_rule is None:
            continue
        # gradFill 제거 + solidFill (yaml이 명시한 색) 삽입
        grad = sp_pr.find(f"{{{A_NS}}}gradFill")
        if grad is not None:
            idx = list(sp_pr).index(grad)
            sp_pr.remove(grad)
        else:
            idx = len(list(sp_pr))
        solid = _etree.SubElement(sp_pr, f"{{{A_NS}}}solidFill")
        sc = matched_rule.get("solid_color") or {"type": "srgbClr", "val": "FFFFFF"}
        clr_type = sc.get("type", "srgbClr")
        clr_val = sc.get("val", "FFFFFF")
        if clr_type == "schemeClr":
            clr = _etree.SubElement(solid, f"{{{A_NS}}}schemeClr")
        else:
            clr = _etree.SubElement(solid, f"{{{A_NS}}}srgbClr")
        clr.set("val", clr_val)
        sp_pr.remove(solid)
        sp_pr.insert(idx, solid)


def fill_cover(slide, sd):
    # gradient mask solidify는 clear_non_placeholder_text에서 *모든 슬라이드 공통* 처리됨
    contents = sd.get("내용", {})
    사업명 = contents.get("사업명", sd.get("제목", ""))
    발주처 = contents.get("발주처", "")
    제안사 = contents.get("제안사", "")
    제출일 = contents.get("제출일", "")

    # 매칭 추적: 어떤 슬롯이 placeholder로 채워졌는지
    title_filled = False
    sub1_filled = False
    sub2_filled = False

    title = find_shape_by_slot(slide, "표지", "title")
    if title:
        # 견본 폰트 보존 (색·굵기·이름)
        set_run_text(title, 사업명, keep_font=True)

        # 박스 dimension — auto_fit 모드면 견본 박스 그대로 사용, 아니면 style.yaml 값
        box_w = COVER_TITLE_BOX.get("width", BODY_AREA["width"])
        box_h = COVER_TITLE_BOX.get("height", 2.4)
        title.left = Inches(COVER_TITLE_BOX.get("left", BODY_AREA["left"]))
        title.top = Inches(COVER_TITLE_BOX["top"])
        title.width = Inches(box_w)
        title.height = Inches(box_h)
        # 자연 줄바꿈
        title.text_frame.word_wrap = True

        # 폰트 크기 결정
        if COVER_AUTO_FIT:
            size = fit_font_to_box(사업명, box_w, box_h)
        else:
            size = fit_title_size(사업명)
        for p in title.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(size)
        _force_white_fill(title)
        title_filled = True

    sub1 = find_shape_by_slot(slide, "표지", "sub1")
    if sub1:
        s1_parts = [p for p in [발주처, 제안사] if p]
        if s1_parts:
            try:
                sub1.left = Inches(COVER_SUB1_BOX.get("left", BODY_AREA["left"]))
                sub1.width = Inches(COVER_SUB1_BOX.get("width", BODY_AREA["width"]))
                if "top" in COVER_SUB1_BOX:
                    sub1.top = Inches(COVER_SUB1_BOX["top"])
                if "height" in COVER_SUB1_BOX:
                    sub1.height = Inches(COVER_SUB1_BOX["height"])
            except Exception:
                pass
            _cover_sub_fit(
                sub1, " / ".join(s1_parts), slide,
                body_left=BODY_AREA["left"], body_w=BODY_AREA["width"],
                max_pt=COVER_TYPOGRAPHY.get("title_max_pt", 22),
            )
            _force_white_fill(sub1)
            sub1_filled = True
    sub2 = find_shape_by_slot(slide, "표지", "sub2")
    if sub2:
        if 제출일:
            try:
                sub2.left = Inches(COVER_SUB2_BOX.get("left", BODY_AREA["left"]))
                sub2.width = Inches(COVER_SUB2_BOX.get("width", BODY_AREA["width"]))
                if "top" in COVER_SUB2_BOX:
                    sub2.top = Inches(COVER_SUB2_BOX["top"])
                if "height" in COVER_SUB2_BOX:
                    sub2.height = Inches(COVER_SUB2_BOX["height"])
            except Exception:
                pass
            _cover_sub_fit(
                sub2, str(제출일), slide,
                body_left=BODY_AREA["left"], body_w=BODY_AREA["width"],
                max_pt=COVER_TYPOGRAPHY.get("sub_max_pt", 18),
            )
            _force_white_fill(sub2)
            sub2_filled = True

    # Fallback: 견본 placeholder 매칭 실패 시 새 textbox 추가
    if not title_filled and 사업명:
        fb_w = COVER_TITLE_BOX.get("width", BODY_AREA["width"])
        fb_h = COVER_TITLE_BOX.get("height", 2.4)
        fb_pt = fit_font_to_box(사업명, fb_w, fb_h) if COVER_AUTO_FIT else fit_title_size(사업명)
        _add_cover_textbox(
            slide, 사업명,
            {"left": COVER_TITLE_BOX.get("left", BODY_AREA["left"]),
             "top": COVER_TITLE_BOX.get("top", 3.0),
             "width": fb_w, "height": fb_h},
            default_pt=fb_pt, color=NAVY, bold=True
        )
    if not sub1_filled and (발주처 or 제안사):
        s1_parts = [p for p in [발주처, 제안사] if p]
        _add_cover_textbox(
            slide, " / ".join(s1_parts),
            {"left": COVER_SUB1_BOX.get("left", BODY_AREA["left"]),
             "top": COVER_SUB1_BOX.get("top", 5.0),
             "width": COVER_SUB1_BOX.get("width", BODY_AREA["width"]),
             "height": COVER_SUB1_BOX.get("height", 0.5)},
            default_pt=14, color=GREY, bold=False
        )
    if not sub2_filled and 제출일:
        _add_cover_textbox(
            slide, 제출일,
            {"left": COVER_SUB2_BOX.get("left", BODY_AREA["left"]),
             "top": COVER_SUB2_BOX.get("top", 5.7),
             "width": COVER_SUB2_BOX.get("width", BODY_AREA["width"]),
             "height": COVER_SUB2_BOX.get("height", 0.5)},
            default_pt=12, color=GREY, bold=False
        )


def _toc_section_parts(sec, default_num):
    """목차 섹션 항목을 (번호, 제목) 으로 정규화.
    - dict   → (sec['번호'], sec['제목'])
    - str    → '1. 사업이해' 패턴 분리, 못 잡으면 (default, str)
    """
    if isinstance(sec, dict):
        num = sec.get("번호", default_num)
        title = sec.get("제목", "")
        num_str = str(num).strip() if num is not None else default_num
        return num_str, str(title).strip()
    s = str(sec).strip()
    import re as _re
    m = _re.match(r"^(\d+)[\.\)]?\s*(.+)$", s)
    if m:
        num_raw = m.group(1)
        return num_raw.zfill(2), m.group(2).strip()
    return default_num, s


def fill_toc(slide, sd, prs=None):
    contents = sd.get("내용", {})
    sections = contents.get("섹션", [])

    title_filled = False
    sections_filled = False

    title = find_shape_by_slot(slide, "목차", "title")
    if title:
        set_run_text(title, sd.get("제목", DEFAULTS["toc_title"]))
        title_filled = True

    if sections:
        n = len(sections)
        half = (n + 1) // 2
        left_items = sections[:half]
        right_items = sections[half:]

        # dict / str 모두 (번호, 제목) 으로 정규화
        left_parts = [_toc_section_parts(s, str(i + 1).zfill(2))
                      for i, s in enumerate(left_items)]
        right_parts = [_toc_section_parts(s, str(i + 1 + half).zfill(2))
                       for i, s in enumerate(right_items)]

        left_nums = "\n".join(num for num, _ in left_parts)
        right_nums = "\n".join(num for num, _ in right_parts)

        for slot_name, content in [("nums_left", left_nums), ("nums_right", right_nums)]:
            s = find_shape_by_slot(slide, "목차", slot_name)
            if s and content:
                set_run_text(s, content)
                sections_filled = True

        left_titles = "\n".join(title for _, title in left_parts)
        right_titles = "\n".join(title for _, title in right_parts)
        for slot_name, content in [("titles_left", left_titles), ("titles_right", right_titles)]:
            s = find_shape_by_slot(slide, "목차", slot_name)
            if s and content:
                set_run_text(s, content)
                sections_filled = True

    # Fallback: 견본 매칭 모두 실패 시 코드 fallback (draw_toc_slide)
    if not (title_filled or sections_filled):
        if prs is None:
            try:
                prs = slide.part.package.presentation_part.presentation
            except Exception:
                pass
        if prs is not None:
            draw_toc_slide(slide, sd, prs)


def fill_section_divider(slide, sd):
    """간지 슬라이드 채움. *통합 title*과 *별도 chap_no slot* 두 모드.

    style.yaml의 section_divider.unified_title_format이 정의되면 *통합 모드*:
    title slot 하나에 포맷된 통합 텍스트 (예: "02.\n시장 현황") 박음.
    그렇지 않으면 *분리 모드*: title slot + chap_no slot 각각.

    style.yaml의 section_divider.draw_chap_no_fallback이 false면 chap_no fallback box 안 그림.
    """
    chap_name = sd.get("장이름") or sd.get("제목", "")
    sn = sd.get("장번호") or sd.get("번호", "")
    try:
        chap_no_int = int(sn) if sn not in (None, "") else 0
    except (TypeError, ValueError):
        chap_no_int = 0
    num_text = f"{chap_no_int:02d}" if chap_no_int else str(sn)

    unified_fmt = SECTION_DIVIDER_CFG.get("unified_title_format", "")
    # 3가지 형식 지원 (yaml의 unified_title_format 값 형태로 분기):
    #   1) 문자열 "{chap_no:02d}.\n{chap_name}" → 단일 title slot에 통합 텍스트
    #   2) dict {chap_no_box: {position,font_pt,color}, title_box: {...}} → *별도 박스 둘* 자동 그림
    #   3) 빈 값 → 분리 모드 (title slot + chap_no slot 각각)
    title_text = chap_name
    use_unified_boxes = isinstance(unified_fmt, dict)

    if isinstance(unified_fmt, str) and unified_fmt:
        title_text = unified_fmt.format(chap_no=chap_no_int, chap_name=chap_name)

    title_filled = False
    num_filled = False

    if use_unified_boxes:
        # 각 박스를 yaml 명시 좌표·폰트로 그림. PY는 단순 데이터 lookup.
        cfg_no = unified_fmt.get("chap_no_box", {})
        cfg_title = unified_fmt.get("title_box", {})
        for cfg, text, default_color in [(cfg_no, num_text, NAVY),
                                          (cfg_title, chap_name, NAVY)]:
            if not cfg or not text:
                continue
            pos = cfg.get("position", {})
            color = _color_by_name(cfg.get("color", "navy"), default=default_color)
            box = slide.shapes.add_textbox(
                Inches(pos.get("left", 0)), Inches(pos.get("top", 0)),
                Inches(pos.get("width", 2)), Inches(pos.get("height", 1)),
            )
            tf = box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = str(text)
            set_run_font(r, size=cfg.get("font_pt", 36),
                         bold=cfg.get("bold", True), color=color)
        title_filled = True
        num_filled = True
    else:
        title = find_shape_by_slot(slide, "간지", "title")
        if title and title_text:
            set_run_text(title, title_text)
            title_filled = True
        # 통합 모드 (문자열)면 chap_no slot 별도 안 채움
        if not unified_fmt:
            num = find_shape_by_slot(slide, "간지", "chap_no")
            if num and num_text:
                set_run_text(num, num_text)
            num_filled = True

    # Fallback: 견본 매칭 실패 시 새 textbox.
    if not title_filled and title_text:
        b = SECTION_DIVIDER_CFG.get("fallback_title_box", {})
        box = slide.shapes.add_textbox(
            Inches(b.get("left", 1.05)), Inches(b.get("top", 3.5)),
            Inches(b.get("width", 11.20)), Inches(b.get("height", 1.5))
        )
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = title_text
        set_run_font(r, size=SECTION_DIVIDER_CFG.get("fallback_title_font_pt", 36), bold=True, color=NAVY)
    # chap_no fallback — yaml의 draw_chap_no_fallback이 true일 때만, 통합 모드 아닐 때만
    if (not num_filled and num_text and not unified_fmt
            and SECTION_DIVIDER_CFG.get("draw_chap_no_fallback", True)):
        b = SECTION_DIVIDER_CFG.get("fallback_number_box", {})
        box = slide.shapes.add_textbox(
            Inches(b.get("left", 1.05)), Inches(b.get("top", 2.5)),
            Inches(b.get("width", 11.20)), Inches(b.get("height", 1.0))
        )
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = num_text
        set_run_font(r, size=SECTION_DIVIDER_CFG.get("fallback_number_font_pt", 48), bold=True, color=GOLD)


# ─────────────────────────────────────────
# 레이아웃 vocabulary
# ─────────────────────────────────────────
# yaml의 `레이아웃` 키로 슬라이드별 명시적 배치 선택.
# 각 함수는 동일 시그니처 (slide, sd, x, y, w, h) — body 영역 내에서 그림.
# 휴리스틱 없이 *명시된 의도대로* 단순 렌더.

def _layout_text_only(slide, sd, x, y, w, h):
    """텍스트만 (시각요소 무시). 불릿·콘텐츠 카드용."""
    messages = sd.get("핵심메시지", []) or []
    contents_dict = sd.get("내용", {}) if isinstance(sd.get("내용"), dict) else {}
    draw_messages_box(slide, messages, contents_dict, x, y, w, h)


def _first_image(sd):
    """sd의 시각요소 중 첫 image dict 반환 (없으면 None)."""
    for v in (sd.get("시각요소") or []):
        if isinstance(v, dict) and v.get("type") == "image":
            return v
    return None


def _images(sd, max_n=None):
    out = []
    for v in (sd.get("시각요소") or []):
        if isinstance(v, dict) and v.get("type") == "image":
            out.append(v)
            if max_n and len(out) >= max_n:
                break
    return out


def _lp(layout_name, key, default):
    """LAYOUT_PARAMS에서 값 로드. 없으면 default."""
    return LAYOUT_PARAMS.get(layout_name, {}).get(key, default)


def _layout_image_focus(slide, sd, x, y, w, h):
    """이미지 1장 대형 + 상단 1~2줄 짧은 설명."""
    messages = sd.get("핵심메시지", []) or []
    msg_h = _lp("image-focus", "message_height", 0.75) if messages else 0
    gap = _lp("image-focus", "image_gap", 0.15)
    if msg_h:
        draw_messages_box(slide, messages[:2], None, x, y, w, msg_h)
    img_y = y + msg_h + (gap if msg_h else 0)
    img_h = h - msg_h - (gap if msg_h else 0)
    img = _first_image(sd)
    if img:
        draw_image(slide, img.get("path", ""), x, img_y, w, img_h,
                   fit=img.get("fit", "contain"), caption=img.get("caption"))


def _layout_split_image_text(slide, sd, x, y, w, h):
    """좌 이미지 + 우 텍스트. 비율은 layout_vocabulary.yaml."""
    gap = _lp("split-image-text", "gap", 0.25)
    img_r = _lp("split-image-text", "image_width_ratio", 0.55)
    txt_r = _lp("split-image-text", "text_width_ratio", 0.45)
    img_w = (w - gap) * img_r
    txt_w = (w - gap) * txt_r
    txt_x = x + img_w + gap
    img = _first_image(sd)
    if img:
        draw_image(slide, img.get("path", ""), x, y, img_w, h,
                   fit=img.get("fit", "contain"), caption=img.get("caption"))
    messages = sd.get("핵심메시지", []) or []
    contents = sd.get("내용", {}) if isinstance(sd.get("내용"), dict) else {}
    draw_messages_box(slide, messages, contents, txt_x, y, txt_w, h)


def _layout_split_text_image(slide, sd, x, y, w, h):
    """좌 텍스트 + 우 이미지."""
    gap = _lp("split-text-image", "gap", 0.25)
    txt_r = _lp("split-text-image", "text_width_ratio", 0.45)
    img_r = _lp("split-text-image", "image_width_ratio", 0.55)
    txt_w = (w - gap) * txt_r
    img_w = (w - gap) * img_r
    img_x = x + txt_w + gap
    img = _first_image(sd)
    if img:
        draw_image(slide, img.get("path", ""), img_x, y, img_w, h,
                   fit=img.get("fit", "contain"), caption=img.get("caption"))
    messages = sd.get("핵심메시지", []) or []
    contents = sd.get("내용", {}) if isinstance(sd.get("내용"), dict) else {}
    draw_messages_box(slide, messages, contents, x, y, txt_w, h)


def _layout_image_grid_2(slide, sd, x, y, w, h):
    """이미지 2장 가로 비교 + 상단 짧은 메시지."""
    messages = sd.get("핵심메시지", []) or []
    msg_h = _lp("image-grid-2", "message_height", 0.65) if messages else 0
    img_gap_v = _lp("image-grid-2", "image_gap", 0.15)
    grid_gap = _lp("image-grid-2", "grid_gap", 0.25)
    if msg_h:
        draw_messages_box(slide, messages[:2], None, x, y, w, msg_h)
    img_y = y + msg_h + (img_gap_v if msg_h else 0)
    img_h = h - msg_h - (img_gap_v if msg_h else 0)
    imgs = _images(sd, max_n=2)
    if len(imgs) >= 2:
        each_w = (w - grid_gap) / 2
        for i, v in enumerate(imgs):
            draw_image(slide, v.get("path", ""), x + i * (each_w + grid_gap),
                       img_y, each_w, img_h,
                       fit=v.get("fit", "contain"), caption=v.get("caption"))
    elif len(imgs) == 1:
        v = imgs[0]
        draw_image(slide, v.get("path", ""), x, img_y, w, img_h,
                   fit=v.get("fit", "contain"), caption=v.get("caption"))


def _layout_visual_only(slide, sd, x, y, w, h):
    """시각요소 1개를 본문 영역 거의 전체에 + 상단 짧은 메시지."""
    messages = sd.get("핵심메시지", []) or []
    msg_h = _lp("visual-only", "message_height", 0.65) if messages else 0
    if msg_h:
        draw_messages_box(slide, messages[:2], None, x, y, w, msg_h)
    vgap = _lp("visual-only", "visual_gap", 0.15)
    v_y = y + msg_h + (vgap if msg_h else 0)
    v_h = h - msg_h - (vgap if msg_h else 0)
    visuals = [v for v in (sd.get("시각요소") or []) if isinstance(v, dict)]
    if visuals:
        v = visuals[0]
        fn = VISUAL_DISPATCH.get(v.get("type"))
        if fn:
            fn(slide, v, x, v_y, w, v_h)


LAYOUT_RENDERERS = {
    "text-only":         _layout_text_only,
    "image-focus":       _layout_image_focus,
    "split-image-text":  _layout_split_image_text,
    "split-text-image":  _layout_split_text_image,
    "image-grid-2":      _layout_image_grid_2,
    "visual-only":       _layout_visual_only,
}


def fill_content(slide, sd):
    """본문: 상단 제목/챕터, 하단 본문(메시지/시각요소).

    `레이아웃` 키가 명시되면 해당 layout renderer로 dispatch.
    명시 안 되거나 모르는 값이면 기존 휴리스틱 (시각요소 + 메시지 자동 배치).
    """
    title = find_shape_by_slot(slide, "본문", "title")
    if title:
        set_run_text(title, sd.get("제목", ""))

    chap_no = find_shape_by_slot(slide, "본문", "chap_no")
    if chap_no:
        n = sd.get("번호", "")
        n_str = f"{int(n):02d}" if str(n).isdigit() else str(n)
        # 모든 placeholder 공통 fit 원칙 적용 (판단 없음 — 텍스트가 박스에 맞도록 폰트 자동 축소)
        _placeholder_fit_text(chap_no, n_str,
                              max_fallback_pt=PLACEHOLDER_FIT_MAX_PT.get("chap_no", 28),
                              slide=slide)
        try:
            chap_no.text_frame.word_wrap = False
        except Exception:
            pass

    chap_name = find_shape_by_slot(slide, "본문", "chap_name")
    if chap_name:
        # yaml의 '종류' 값 그대로 사용. 종류 없으면 style.yaml의 default
        set_run_text(chap_name, sd.get("종류", DEFAULTS["chap_name_body"]))
        try:
            chap_name.text_frame.word_wrap = False
        except Exception:
            pass

    # 견본의 더미 본문 textbox 식별 (style.yaml의 anchor 사용)
    body_shape = find_shape_by_pos(
        slide, BODY_TEXTBOX_ANCHOR["left"], BODY_TEXTBOX_ANCHOR["top"],
        tol=SHAPE_FINDING.get("body_tolerance", 0.5)
    )
    if body_shape is None:
        body_shape = find_shape_by_text(slide, "텍스트 서술형 내용")

    # 시각요소 구조화 검출
    visuals = sd.get("시각요소", [])
    structured_visuals = [v for v in visuals if isinstance(v, dict) and "type" in v]
    messages = sd.get("핵심메시지", [])
    contents_dict = sd.get("내용", {}) if isinstance(sd.get("내용"), dict) else {}

    # 본문 영역 (style.yaml에서 로드)
    body_left = BODY_AREA["left"]
    body_top = BODY_AREA["top"]
    body_width = BODY_AREA["width"]
    body_height = BODY_AREA["height"]

    # title shape 충돌 회피 — title.bottom + gap 보다 body_top이 위면 자동 조정.
    # 정책 (gap 값)은 system_defaults.yaml content_layout.title_gap.
    if title and title.top is not None and title.height is not None:
        title_bottom_in = (title.top + title.height) / 914400
        title_gap = CONTENT_LAYOUT_CFG.get("title_gap", 0.15)
        if title_bottom_in + title_gap > body_top:
            new_body_top = title_bottom_in + title_gap
            body_height = max(1.0, body_height - (new_body_top - body_top))
            body_top = new_body_top

    # ─── 견본의 더미 라벨·빈 박스 정리 (시각요소 그리기 전에) ───
    # placeholder 정리는 main loop에서 일괄 처리됨 (모든 placeholder text 비움).
    # 여기는 비-placeholder textbox만 처리.
    dummy_keywords = SAMPLE_CLONE.get("dummy_texts") or []
    keep_ids = {id(s) for s in (body_shape, title, chap_no, chap_name) if s is not None}
    for shape in list(slide.shapes):
        if id(shape) in keep_ids:
            continue
        if shape.is_placeholder:
            continue
        if shape.has_text_frame:
            txt = shape.text_frame.text.strip()
            if any(kw in txt for kw in dummy_keywords):
                sp = shape._element
                sp.getparent().remove(sp)
                continue
            if not txt:
                sp = shape._element
                sp.getparent().remove(sp)
                continue
        else:
            if shape.width and shape.height:
                w_in = shape.width / 914400
                h_in = shape.height / 914400
                if w_in > 0.5 and h_in > 0.5:
                    sp = shape._element
                    sp.getparent().remove(sp)
                    continue

    # ─── 명시 레이아웃 dispatch (휴리스틱 우회) ───
    layout_name = sd.get("레이아웃")
    if layout_name in LAYOUT_RENDERERS:
        keep_shapes_check = [s for s in (title, chap_no, chap_name) if s is not None]
        if body_shape and not any(shapes_same(body_shape, ks) for ks in keep_shapes_check):
            set_run_text(body_shape, "")
        LAYOUT_RENDERERS[layout_name](
            slide, sd, body_left, body_top, body_width, body_height
        )
        return

    # ─── 명시 레이아웃 없음 — 휴리스틱 fallback ───
    if structured_visuals:
        # 분할 모드: 메시지(상단) + 시각요소(하단)
        keep_shapes_check = [s for s in (title, chap_no, chap_name) if s is not None]
        if body_shape and not any(shapes_same(body_shape, ks) for ks in keep_shapes_check):
            set_run_text(body_shape, "")

        # 메시지 영역 높이 — 메시지 수 + 이미지 위주 슬라이드 여부에 따라 가변
        msg_count = len(messages or [])
        content_count = len(contents_dict) if contents_dict else 0
        total_msg = msg_count + content_count
        img_count = sum(1 for v in structured_visuals if v.get("type") == "image")
        # 메시지 영역 높이 공식: system_defaults.yaml content_layout 섹션
        cl_with = CONTENT_LAYOUT_CFG.get("with_visuals", {})
        cl_no = CONTENT_LAYOUT_CFG.get("no_visuals", {})
        if total_msg == 0:
            msg_h = 0.0
        elif img_count >= 1:
            msg_h = min(cl_with.get("cap_inch", 1.4),
                        cl_with.get("base_inch", 0.4) + cl_with.get("per_line_inch", 0.25) * total_msg)
        else:
            msg_h = min(cl_no.get("cap_inch", 2.2),
                        cl_no.get("base_inch", 0.4) + cl_no.get("per_line_inch", 0.4) * total_msg)

        if msg_h > 0:
            draw_messages_box(slide, messages, contents_dict,
                              body_left, body_top, body_width, msg_h)
        viz_top = body_top + msg_h + (0.2 if msg_h > 0 else 0)
        viz_h = body_height - msg_h - (0.2 if msg_h > 0 else 0)
        render_visuals(slide, structured_visuals,
                       body_left, viz_top, body_width, viz_h)
    else:
        # 시각요소 없는 본문 — 항상 *새 textbox 추가* (견본 placeholder 의존 X)
        # 견본 body_shape 매칭 실패해도 콘텐츠 누락 없음
        lines = []
        for m in messages:
            lines.append(f"▶  {m}")
        if contents_dict:
            if lines:
                lines.append("")
            for k, v in contents_dict.items():
                if isinstance(v, list):
                    lines.append(f"■ {k}:")
                    for sub in v:
                        lines.append(f"    · {sub}")
                else:
                    lines.append(f"■ {k}: {v}")
        if visuals:
            if lines:
                lines.append("")
            lines.append("[시각요소]")
            for v in visuals:
                lines.append(f"  · {v}")

        if not lines:
            return  # 그릴 콘텐츠 없음

        n_lines = max(len(lines), 1)
        # 라인 수 기반 폰트 선택 (system_defaults.yaml body_text_font_by_lines)
        msg_font = 11
        for rule in BODY_TEXT_FONT_BY_LINES:
            if n_lines <= rule.get("max_lines", 9999):
                msg_font = rule.get("font_pt", 11)
                break

        # 견본의 본문 textbox는 (있어도) 비우고 *새 textbox*에 콘텐츠 그림
        # → 어떤 견본이든 콘텐츠 누락 없음
        if body_shape:
            keep_shapes_check = [s for s in (title, chap_no, chap_name) if s is not None]
            if not any(shapes_same(body_shape, ks) for ks in keep_shapes_check):
                set_run_text(body_shape, "")  # 견본 더미 비움

        new_box = slide.shapes.add_textbox(
            Inches(body_left), Inches(body_top),
            Inches(body_width), Inches(body_height)
        )
        tf = new_box.text_frame
        tf.word_wrap = True
        first = True
        for line in lines:
            if first:
                p = tf.paragraphs[0]
                first = False
            else:
                p = tf.add_paragraph()
            r = p.add_run()
            r.text = line
            set_run_font(r, size=msg_font, bold=True, color=DARK)
            p.space_after = Pt(12)


def fill_company_intro(slide, sd):
    headline_filled = False
    sub_filled = False

    title = find_shape_by_slot(slide, "본문", "title")
    if title:
        set_run_text(title, sd.get("제목", ""))
    chap_no = find_shape_by_slot(slide, "본문", "chap_no")
    if chap_no:
        n = sd.get("번호", "")
        set_run_text(chap_no, f"{int(n):02d}" if str(n).isdigit() else str(n))
    chap_name = find_shape_by_slot(slide, "본문", "chap_name")
    if chap_name:
        # yaml의 '종류' 값 그대로 사용. 종류 없으면 style.yaml의 default
        set_run_text(chap_name, sd.get("종류", DEFAULTS["chap_name_company_intro"]))
    ha = COMPANY_INTRO_ANCHORS["headline"]
    headline = find_shape_by_pos(slide, ha["left"], ha["top"], tol=SHAPE_FINDING.get("company_intro_tolerance", 0.4))
    if headline:
        msgs = sd.get("핵심메시지", [])
        if msgs:
            set_run_text(headline, msgs[0])
            headline_filled = True
    sa = COMPANY_INTRO_ANCHORS["sub"]
    sub = find_shape_by_pos(slide, sa["left"], sa["top"], tol=SHAPE_FINDING.get("company_intro_tolerance", 0.4))
    if sub:
        msgs = sd.get("핵심메시지", [])
        sub_lines = msgs[1:] if len(msgs) > 1 else []
        contents = sd.get("내용", {})
        if isinstance(contents, dict):
            for k, v in contents.items():
                if isinstance(v, list):
                    sub_lines.append(f"{k}: {', '.join(str(x) for x in v)}")
                else:
                    sub_lines.append(f"{k}: {v}")
        if sub_lines:
            set_run_text(sub, "\n".join(sub_lines))
            sub_filled = True
    # (회사 식별 텍스트 제거는 sample_clone.dummy_texts (style.yaml) 가 자동 처리)

    # Fallback: headline·sub 모두 매칭 실패 시 본문 모드로 처리 (콘텐츠 누락 방지)
    if not (headline_filled or sub_filled):
        fill_content(slide, sd)


def _shape_text_anchor(shape):
    """텍스트 박스의 anchor (top/center/bottom)을 반환. 못 읽으면 None."""
    try:
        body_pr = shape.text_frame._txBody.find(
            "{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr"
        )
        if body_pr is not None:
            return body_pr.get("anchor")  # 't'/'ctr'/'b'
    except Exception:
        pass
    return None


def _safe_height_for_placeholder(slide, shape):
    """slide의 다른 shape들이 placeholder 영역에 침범하는지 확인하고
    텍스트 anchor 방향으로 안전한 height 반환. 침범 없으면 원래 height."""
    if not shape.width or not shape.height:
        return 0
    ph_x = shape.left / 914400
    ph_y = shape.top / 914400
    ph_w = shape.width / 914400
    ph_h = shape.height / 914400
    ph_mid = ph_y + ph_h / 2
    ph_bottom = ph_y + ph_h

    anchor = _shape_text_anchor(shape) or "t"

    safe_top = ph_y
    safe_bottom = ph_bottom

    for other in slide.shapes:
        if other._element is shape._element:
            continue
        if not other.width or not other.height:
            continue
        o_x = (other.left or 0) / 914400
        o_y = (other.top or 0) / 914400
        o_w = (other.width or 0) / 914400
        o_h = (other.height or 0) / 914400
        # 가로 겹침 (system_defaults.yaml shape_safety.overlap_threshold)
        overlap_w = max(0, min(ph_x + ph_w, o_x + o_w) - max(ph_x, o_x))
        if overlap_w < min(ph_w, o_w) * SHAPE_SAFETY.get("overlap_threshold", 0.2):
            continue
        # 완전히 위·아래는 무관
        if o_y + o_h <= ph_y or o_y >= ph_bottom:
            continue
        # 침범 위치 분류
        o_mid = o_y + o_h / 2
        if o_mid < ph_mid:
            # 위쪽 침범
            if anchor in ("b", "ctr"):
                safe_top = max(safe_top, o_y + o_h)
        else:
            # 아래쪽 침범
            if anchor in ("t", "ctr"):
                safe_bottom = min(safe_bottom, o_y)

    return max(SHAPE_SAFETY.get("min_safe_height", 0.3), safe_bottom - safe_top)


def _max_safe_width_right(slide, shape, body_left, body_w):
    """shape의 우측으로 안전하게 확장 가능한 max width 계산.
    세로로 겹치는 다른 shape 회피. body_area의 right edge 이내로 제한."""
    box_left = (shape.left / 914400) if shape.left else body_left
    box_top = (shape.top / 914400) if shape.top else 0
    box_h = (shape.height / 914400) if shape.height else 0
    box_bottom = box_top + box_h
    right_limit = body_left + body_w
    for other in slide.shapes:
        if other._element is shape._element:
            continue
        if not other.width or not other.height:
            continue
        o_x = (other.left or 0) / 914400
        o_y = (other.top or 0) / 914400
        o_w = (other.width or 0) / 914400
        o_h = (other.height or 0) / 914400
        # 세로로 겹치지 않으면 무관
        if o_y + o_h <= box_top or o_y >= box_bottom:
            continue
        # box 오른쪽에 있고 right_limit보다 가까운 shape이면 right_limit 갱신
        if o_x > box_left + 0.5 and o_x < right_limit:
            right_limit = max(box_left + 1.0, o_x - 0.15)
    return max(SHAPE_SAFETY.get("min_safe_width", 0.5), right_limit - box_left)


def _cover_sub_fit(shape, text, slide, body_left, body_w,
                   max_pt=None, min_pt=None, min_readable_pt=None):
    """표지 sub 박스에 텍스트 fit. 폰트가 min_readable 미만이 되면
    우측 adjacency-safe 한계까지 박스 width 확장 후 재 fit.
    기본 폰트 범위는 system_defaults.yaml cover_typography 섹션."""
    if max_pt is None:
        max_pt = COVER_TYPOGRAPHY.get("sub_fit_max_pt", 22)
    if min_pt is None:
        min_pt = COVER_TYPOGRAPHY.get("sub_fit_min_pt", 12)
    if min_readable_pt is None:
        min_readable_pt = COVER_TYPOGRAPHY.get("sub_fit_min_readable_pt", 14)
    if not shape or not shape.has_text_frame:
        return
    set_run_text(shape, text, keep_font=True)
    shape.text_frame.word_wrap = True
    box_w = (shape.width / 914400) if shape.width else 0
    box_h = (shape.height / 914400) if shape.height else 0
    if box_w <= 0 or box_h <= 0:
        return
    safe_h = _safe_height_for_placeholder(slide, shape) if slide else box_h
    font_pt = fit_font_to_box(text, box_w, safe_h, max_pt=max_pt, min_pt=min_pt)
    if font_pt < min_readable_pt and slide is not None:
        max_w = _max_safe_width_right(slide, shape, body_left, body_w)
        width_expand_thr = SHAPE_SAFETY.get("width_expand_threshold", 1.0)
        if max_w > box_w + width_expand_thr:
            shape.width = Inches(max_w)
            box_w = max_w
            font_pt = fit_font_to_box(text, box_w, safe_h, max_pt=max_pt, min_pt=min_pt)
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            r.font.size = Pt(font_pt)


def _placeholder_fit_text(shape, text, max_fallback_pt=None, slide=None):
    """placeholder shape에 텍스트를 fit. 원래 폰트 보존하되 박스 + 인접 shape를
    고려한 안전 영역 기반으로 폰트 크기 자동 축소.
    max_fallback_pt 미지정 시 system_defaults.yaml placeholder_fit_max_pt.default.
    """
    if max_fallback_pt is None:
        max_fallback_pt = PLACEHOLDER_FIT_MAX_PT.get("default", 40)
    if not shape.has_text_frame:
        return
    # 원래 폰트 크기 추출
    original_pt = None
    try:
        tf0 = shape.text_frame
        if tf0.paragraphs and tf0.paragraphs[0].runs:
            sz = tf0.paragraphs[0].runs[0].font.size
            if sz is not None:
                original_pt = sz.pt
    except Exception:
        pass

    set_run_text(shape, text, keep_font=True)
    shape.text_frame.word_wrap = True

    box_w = (shape.width / 914400) if shape.width else 0
    box_h = (shape.height / 914400) if shape.height else 0
    if box_w <= 0 or box_h <= 0:
        return

    # 인접 shape 침범 고려한 안전 height
    safe_h = box_h
    if slide is not None:
        safe_h = _safe_height_for_placeholder(slide, shape)

    upper = int(original_pt) if original_pt else max_fallback_pt
    target = fit_font_to_box(text, box_w, safe_h, max_pt=upper, min_pt=12)
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            r.font.size = Pt(target)


def fill_thanks(slide, sd):
    msg = sd.get("내용", {}).get("메시지", sd.get("제목", DEFAULTS["thanks_message"]))
    title_filled = False
    title = find_shape_by_slot(slide, "감사", "title")
    if title and msg:
        _placeholder_fit_text(title, msg,
                              max_fallback_pt=PLACEHOLDER_FIT_MAX_PT.get("thanks_title", 60),
                              slide=slide)
        title_filled = True

    # Fallback. 박스·폰트는 system_defaults.yaml thanks 섹션.
    if not title_filled and msg:
        b = THANKS_CFG.get("fallback_box", {})
        bl = b.get("left", 0.5); bt = b.get("top", 3.0)
        bw = b.get("width", 12.33); bh = b.get("height", 1.5)
        box = slide.shapes.add_textbox(Inches(bl), Inches(bt), Inches(bw), Inches(bh))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = msg
        size = fit_font_to_box(msg, bw, bh, max_pt=THANKS_CFG.get("fallback_max_font_pt", 40))
        set_run_font(r, size=size, bold=True, color=NAVY)


def _draw_label_overlay(slide, cfg, text, align):
    """layout에 박힌 라벨을 *흰 박스로 가린 뒤 우리 텍스트로 덮어씀*.

    박스 *width 자동 fit*: cfg.position 에 left·width 가 명시되어 있으면 그대로 사용.
    right + max_width 가 명시되어 있으면 *텍스트 길이 기반*으로 width 계산:
    - 텍스트 폭 추정: char_count * font_pt * char_width_factor / 72 (inch)
    - 추정 폭이 max_width보다 작으면 폭만큼 사용, 박스 right edge는 cfg.position.right.
    - 추정 폭이 max_width보다 크면 font_pt를 비례 축소해 max_width에 맞춤.

    cfg는 style.yaml의 top_right_label·left_chap_no 같은 dict.
    align: 'right' (우상단) | 'left' | 'center' (좌상단 chap_no는 center).
    """
    if not cfg.get("enabled") or text is None or text == "":
        return
    pos = cfg.get("position", {})
    color = _color_by_name(cfg.get("color", "navy"))
    font_pt = cfg.get("font_pt", 14)
    char_w_factor = TYPOGRAPHY_FACTORS.get("char_width", 0.9)
    pad = TOP_RIGHT_LABEL_PADDING
    height = pos.get("height", 0.35)

    text_str = str(text)
    # 텍스트 추정 폭 (inch) — 한글 ≈ 1 char, 영문/숫자 ≈ 0.6 char (단순화)
    char_units = sum(1.0 if ord(c) > 127 else 0.6 for c in text_str)
    estimated_w = char_units * font_pt * char_w_factor / 72

    if "right" in pos:
        # right anchor 모드: width 자동
        max_w = pos.get("max_width", pos.get("width", 4.0))
        # 폰트 축소 시 최소 pt — yaml의 fit_font_box.min_pt 재사용
        min_font_pt = FIT_FONT_BOX.get("min_pt", 8)
        if estimated_w > max_w:
            scale = max_w / estimated_w
            font_pt = max(min_font_pt, int(font_pt * scale))
            estimated_w = max_w
        # margin 여유분 (양쪽 textbox margin + pad)
        total_margin = (TEXTBOX_MARGIN.get("left", 0.05)
                        + TEXTBOX_MARGIN.get("right", 0.05) + 0.1)
        box_w = min(estimated_w + total_margin, max_w)
        left = max(0, pos["right"] - box_w)
    else:
        # left + width 명시 모드 (기존)
        left = pos.get("left", 0)
        box_w = pos.get("width", 2.1)

    box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(max(0, left - pad)),
        Inches(max(0, pos.get("top", 0) - pad)),
        Inches(box_w + 2 * pad),
        Inches(height + 2 * pad),
    )
    box.fill.solid()
    box.fill.fore_color.rgb = WHITE
    box.line.fill.background()

    tf = box.text_frame
    tf.margin_left = Inches(TEXTBOX_MARGIN.get("left", 0.05))
    tf.margin_right = Inches(TEXTBOX_MARGIN.get("right", 0.05))
    tf.margin_top = Inches(TEXTBOX_MARGIN.get("top", 0.02))
    tf.margin_bottom = Inches(TEXTBOX_MARGIN.get("bottom", 0.02))
    p = tf.paragraphs[0]
    align_map = {"right": PP_ALIGN.RIGHT, "left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER}
    p.alignment = align_map.get(align, PP_ALIGN.LEFT)
    r = p.add_run()
    r.text = text_str
    set_run_font(r, size=font_pt, bold=True, color=color)


def _format_label(fmt, chap_no, chap_name):
    """style.yaml의 format 문자열로 라벨 텍스트 만듦.
    chap_no는 정수로 변환 시도, 실패하면 0. chap_name 빈 채면 빈 문자열.
    format 자체가 비면 빈 문자열 (정책 없음)."""
    if not fmt:
        return ""
    try:
        n = int(chap_no) if chap_no not in (None, "") else 0
    except (TypeError, ValueError):
        n = 0
    return fmt.format(chap_no=n, chap_name=chap_name or "")


def draw_chapter_labels(slide, sd, kind=None):
    """yaml의 chapter_labels list 따라 라벨 N개를 layout 위에 덮어 그림.
    각 라벨의 enabled_for_kinds (yaml) 에 *kind*가 포함된 경우만 그림.
    enabled_for_kinds 미지정이면 모든 kind에 그림.
    """
    chap_no = sd.get("장번호")
    chap_name = sd.get("장이름")
    for cfg in CHAPTER_LABELS:
        if not cfg.get("enabled"):
            continue
        kinds = cfg.get("enabled_for_kinds")
        if kinds and kind and kind not in kinds:
            continue
        fmt = cfg.get("format")
        if not fmt:
            continue
        text = _format_label(fmt, chap_no, chap_name)
        align = cfg.get("align", "left")
        _draw_label_overlay(slide, cfg, text, align=align)


# ─────────────────────────────────────────
# 코드 fallback — template에 견본 없을 때
# ─────────────────────────────────────────

def make_blank_slide(prs):
    """가장 빈 레이아웃으로 빈 슬라이드 추가. placeholder 모두 제거."""
    master = prs.slide_masters[0]
    layouts = list(master.slide_layouts)
    # placeholder 수가 가장 적은 레이아웃 선택
    best = min(layouts, key=lambda L: len(list(L.placeholders)))
    new_slide = prs.slides.add_slide(best)
    # placeholder 제거 (자리 차지 방지)
    for shp in list(new_slide.shapes):
        if shp.is_placeholder:
            sp = shp._element
            sp.getparent().remove(sp)
    return new_slide


def draw_toc_slide(slide, sd, prs):
    """코드로 목차 직접 그리기. 박스·위치는 system_defaults.yaml toc 섹션."""
    slide_w = prs.slide_width / 914400
    slide_h = prs.slide_height / 914400

    contents = sd.get("내용", {})
    sections = contents.get("섹션", [])

    title_cfg = TOC_CFG.get("title_box", {})
    accent_cfg = TOC_CFG.get("accent_line", {})
    sec_num_cfg = TOC_CFG.get("section_number_box", {})
    sec_title_cfg = TOC_CFG.get("section_title_box", {})

    # 상단 제목
    title_box = slide.shapes.add_textbox(
        Inches(title_cfg.get("left", 0.6)), Inches(title_cfg.get("top", 0.4)),
        Inches(slide_w - title_cfg.get("width_from_right", 1.2)),
        Inches(title_cfg.get("height", 0.7))
    )
    p = title_box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = sd.get("제목", DEFAULTS["toc_title"])
    set_run_font(r, size=TYPOGRAPHY["toc_title_pt"], bold=True, color=NAVY)

    # 골드 강조선 (제목 아래)
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(accent_cfg.get("left", 0.6)), Inches(accent_cfg.get("top", 1.15)),
        Inches(accent_cfg.get("width", 2.0)), Inches(accent_cfg.get("height", 0.05))
    )
    line.fill.solid()
    line.fill.fore_color.rgb = GOLD
    line.line.fill.background()

    if not sections:
        return

    n = len(sections)
    half = (n + 1) // 2
    left_items = sections[:half]
    right_items = sections[half:]

    body_top = TOC_CFG.get("body_top", 1.7)
    body_bottom = slide_h - TOC_CFG.get("body_bottom_from_bottom", 0.7)
    body_h = body_bottom - body_top
    col_w = (slide_w - TOC_CFG.get("column_width_offset", 1.5)) / 2
    sec_num_w = sec_num_cfg.get("width", 0.9)
    sec_num_h = sec_num_cfg.get("height", 0.6)
    sec_title_left_off = sec_title_cfg.get("left_offset", 1.0)
    sec_title_top_off = sec_title_cfg.get("top_offset", 0.08)
    sec_title_h = sec_title_cfg.get("height", 0.55)
    right_x_offset = TOC_CFG.get("right_column_x_offset", 0.3)

    def render_col(items, x_start, num_offset):
        if not items:
            return
        line_h = body_h / max(len(items), 1)
        for i, sec in enumerate(items):
            num_str, title_str = _toc_section_parts(
                sec, str(i + 1 + num_offset).zfill(2)
            )
            y = body_top + i * line_h + (line_h - sec_title_h) / 2
            num_box = slide.shapes.add_textbox(
                Inches(x_start), Inches(y),
                Inches(sec_num_w), Inches(sec_num_h)
            )
            nr = num_box.text_frame.paragraphs[0].add_run()
            nr.text = num_str
            set_run_font(nr, size=TYPOGRAPHY["toc_title_pt"], bold=True, color=GOLD)
            tt = slide.shapes.add_textbox(
                Inches(x_start + sec_title_left_off), Inches(y + sec_title_top_off),
                Inches(col_w - sec_title_left_off), Inches(sec_title_h)
            )
            tr = tt.text_frame.paragraphs[0].add_run()
            tr.text = title_str
            set_run_font(tr, size=TYPOGRAPHY["toc_section_num_pt"], bold=True, color=NAVY)

    render_col(left_items, x_start=title_cfg.get("left", 0.6), num_offset=0)
    render_col(right_items, x_start=title_cfg.get("left", 0.6) + col_w + right_x_offset, num_offset=half)


def add_notes(slide, notes_text):
    if not notes_text:
        return
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = notes_text.strip()


# ─────────────────────────────────────────
# main
# ─────────────────────────────────────────

def _user_error(msg, hint=None, code=1):
    """사용자 친화적 에러 메시지 + 종료. traceback 대신 명확한 안내."""
    print(f"\n[오류] {msg}", file=sys.stderr)
    if hint:
        print(f"  → {hint}", file=sys.stderr)
    sys.exit(code)


def _user_warn(msg):
    print(f"[경고] {msg}", file=sys.stderr)


KNOWN_KINDS = {"표지", "목차", "간지", "회사소개", "본문", "결론", "부록", "감사"}
KNOWN_SLIDE_KEYS = {
    "종류", "제목", "부제", "내용", "시각요소",
    "장번호", "장이름", "메시지", "본문",
}

# 시각요소 스키마는 templates/visual_element_schema.yaml 에서 로드.
# PY는 *알고리즘만*, *데이터(키 이름·필수 여부·별칭)는 yaml*.
_VISUAL_SCHEMA_PATH = Path(__file__).parent.parent / "templates" / "visual_element_schema.yaml"
VISUAL_SCHEMA = {}

def _load_visual_schema():
    global VISUAL_SCHEMA
    try:
        with open(_VISUAL_SCHEMA_PATH, "r", encoding="utf-8") as f:
            doc = yaml.safe_load(f) or {}
        VISUAL_SCHEMA = doc.get("visual_elements", {})
    except FileNotFoundError:
        VISUAL_SCHEMA = {}

_load_visual_schema()


def _alias_hint(v_type, bad_key):
    """schema의 common_mistakes 에서 잘못 사용된 키에 대응되는 *올바른 키*를 찾아 안내."""
    spec = VISUAL_SCHEMA.get(v_type, {})
    mistakes = spec.get("common_mistakes", []) or []
    for line in mistakes:
        if bad_key in str(line):
            return str(line)
    return None


def _check_value_type(val, allowed_types, sub_keys_def, ctx):
    """schema의 value_types/item_types 따라 한 값의 형식 검증.
    allowed_types: ['str', 'dict_label_body', ...]
    sub_keys_def: dict 형식의 sub-keys 정의 ({'dict_label_body': {required: [...], optional: [...]}})
    반환: 매칭된 type 이름 또는 None (어느 형식에도 안 맞음).
    """
    if not allowed_types:
        return None
    for type_name in allowed_types:
        if type_name == "str":
            if isinstance(val, str):
                return "str"
        elif type_name == "dict" or type_name.startswith("dict_"):
            if not isinstance(val, dict):
                continue
            # sub_keys_def에 정의되어 있으면 sub-keys 확인
            sub_spec = (sub_keys_def or {}).get(type_name, {})
            req_subs = set(sub_spec.get("required", []))
            if req_subs and not req_subs.issubset(set(val.keys())):
                continue
            return type_name
        elif type_name == "list":
            if isinstance(val, list):
                return "list"
    return None


def _validate_visual(v, ctx):
    """시각요소 dict 검증. 정의는 *visual_element_schema.yaml*. PY는 단순 룰 평가.
    검증 항목:
    1. type 키 존재 + 정의된 type 여부
    2. 필수 키 존재
    3. 각 키의 *값 형식* (value_types/item_types) 일치
    """
    if not isinstance(v, dict):
        _user_warn(f"{ctx} 시각요소가 dict 아님 — 무시.")
        return
    t = v.get("type")
    if not t:
        _user_warn(f"{ctx} 시각요소에 'type' 없음 — 무시.")
        return
    if t not in VISUAL_SCHEMA:
        _user_warn(
            f"{ctx} 시각요소 type='{t}' 미정의 — 무시. "
            f"지원 (visual_element_schema.yaml): {', '.join(sorted(VISUAL_SCHEMA.keys()))}"
        )
        return
    spec = VISUAL_SCHEMA[t]
    required_spec = spec.get("required") or {}
    required_keys = set(required_spec.keys())
    missing = required_keys - set(v.keys())
    if missing:
        # 잘못된 키 (별칭) 사용 흔적 찾기 — schema의 common_mistakes에서 안내
        extra_keys = set(v.keys()) - required_keys - set((spec.get("optional") or {}).keys()) - {"type"}
        hints = []
        for bk in extra_keys:
            h = _alias_hint(t, bk)
            if h:
                hints.append(f'    "{bk}" → {h}')
        hint_msg = ("\n  올바른 키 안내:\n" + "\n".join(hints)) if hints else ""
        _user_warn(
            f"{ctx} 시각요소(type={t}) 필수 키 누락: {sorted(missing)}. 해당 요소 미표시.{hint_msg}"
        )

    # 각 필수 키의 *값 형식* 검증
    for key_name, key_spec in required_spec.items():
        if key_name not in v:
            continue
        val = v[key_name]
        # value_types (dict 의 각 *값*에 대한 형식 — 예: matrix_2x2.quadrants 의 Q1 값)
        if isinstance(key_spec, dict) and "value_types" in key_spec and isinstance(val, dict):
            allowed = key_spec.get("value_types", [])
            sub_def = key_spec.get("value_sub_keys", {})
            for sub_k, sub_v in val.items():
                matched = _check_value_type(sub_v, allowed, sub_def, f"{ctx} {key_name}.{sub_k}")
                if matched is None:
                    _user_warn(
                        f"{ctx} 시각요소(type={t}) 키 '{key_name}.{sub_k}' 값 형식 잘못됨. "
                        f"허용: {allowed}. 받음: {type(sub_v).__name__} {sub_v!r}. "
                        f"→ ppt-designer가 schema yaml 따라 yaml 다시 작성 필요."
                    )
        # item_types (list 의 각 *항목*에 대한 형식 — 예: callout_cards.items, timeline.items)
        if isinstance(key_spec, dict) and "item_types" in key_spec and isinstance(val, list):
            allowed = key_spec.get("item_types", [])
            sub_def = key_spec.get("item_sub_keys", {})
            for i, item in enumerate(val):
                matched = _check_value_type(item, allowed, sub_def, f"{ctx} {key_name}[{i}]")
                if matched is None:
                    _user_warn(
                        f"{ctx} 시각요소(type={t}) 키 '{key_name}[{i}]' 항목 형식 잘못됨. "
                        f"허용: {allowed}. 받음: {type(item).__name__} {item!r}"
                    )

    if t == "image":
        path = v.get("path")
        if path and not Path(path).exists():
            _user_warn(f"{ctx} image path 미존재: {path} — 자리표시자만 표시.")


def validate_content_yaml(data, yaml_path):
    """콘텐츠 yaml 구조 검증. 치명적 오류는 _user_error로 종료, 경미는 _user_warn."""
    if data is None:
        _user_error(
            f"콘텐츠 yaml이 비어있음: {yaml_path}",
            hint="'프레젠테이션:'·'슬라이드:' 섹션이 있어야 합니다."
        )
    if not isinstance(data, dict):
        _user_error(
            f"콘텐츠 yaml의 최상위가 dict가 아님: {yaml_path}",
            hint="yaml 파일은 키:값 형태로 시작해야 합니다."
        )

    slides = data.get("슬라이드")
    if not isinstance(slides, list) or len(slides) == 0:
        _user_error(
            f"콘텐츠 yaml에 '슬라이드' 리스트가 없거나 비어있음: {yaml_path}",
            hint="최소 1개 이상의 슬라이드를 정의해야 합니다."
        )

    cover_seen = False
    for i, sd in enumerate(slides, start=1):
        ctx = f"슬라이드 {i}"
        if not isinstance(sd, dict):
            _user_error(f"{ctx} 가 dict가 아님 (yaml 들여쓰기 확인).")

        kind = sd.get("종류", "본문")
        if kind not in KNOWN_KINDS:
            _user_warn(
                f"{ctx} 종류 '{kind}' 미정의 — 본문으로 처리. "
                f"지원: {', '.join(sorted(KNOWN_KINDS))}"
            )

        # 알 수 없는 최상위 키 (오타 잡기)
        unknown = set(sd.keys()) - KNOWN_SLIDE_KEYS
        if unknown:
            _user_warn(f"{ctx} 알 수 없는 키 {sorted(unknown)} — 무시. 오타 확인 권장.")

        # 시각요소 검증
        visuals = sd.get("시각요소", [])
        if visuals and not isinstance(visuals, list):
            _user_warn(f"{ctx} '시각요소'가 리스트 아님 — 무시.")
        elif isinstance(visuals, list):
            for j, v in enumerate(visuals):
                _validate_visual(v, f"{ctx} 시각요소 {j+1}")

        # 종류별 필수 필드
        if kind == "표지":
            cover_seen = True
            사업명 = (sd.get("내용", {}) or {}).get("사업명") or sd.get("제목")
            if not 사업명:
                _user_warn(f"{ctx} (표지) 에 사업명/제목 없음 — 빈 표지 생성.")
        elif kind == "간지":
            if not (sd.get("장번호") or sd.get("장이름") or sd.get("제목")):
                _user_warn(f"{ctx} (간지) 에 장번호·장이름·제목 모두 없음.")
        elif kind == "본문":
            if not sd.get("제목"):
                _user_warn(f"{ctx} (본문) 에 제목 없음.")

    if not cover_seen:
        _user_warn("표지 슬라이드가 없습니다. 첫 슬라이드를 표지로 권장.")


def main():
    import argparse
    import subprocess

    ap = argparse.ArgumentParser(
        description="콘텐츠 yaml + 회사 template → .pptx 생성"
    )
    ap.add_argument("yaml", help="콘텐츠 yaml 경로")
    ap.add_argument("pptx", help="출력 .pptx 경로")
    ap.add_argument(
        "--template", "-t",
        help="템플릿 이름 (templates/<name>.pptx 또는 풀 경로). "
             "지정 시 templates/<name>.style.yaml 자동 로드."
    )
    ap.add_argument(
        "--style", "-s",
        help="style.yaml 직접 경로 지정 (--template보다 우선)."
    )
    ap.add_argument(
        "--no-sync", action="store_true",
        help="templates/ 자동 sync 스킵 (기본은 sync 실행)."
    )
    args = ap.parse_args()

    src_yaml = args.yaml
    dst_pptx = args.pptx

    # 입력 yaml 존재 검증
    if not Path(src_yaml).exists():
        _user_error(
            f"콘텐츠 yaml 파일 없음: {src_yaml}",
            hint="경로를 확인하세요. (예: output/20260514/slides_*.yaml)"
        )

    # 1) templates/ 자동 sync — 폐지됨.
    # PY 자동 매핑 휴리스틱 제거에 따라 sync 스크립트도 폐지.
    # 새 template은 /onboard-template 슬래시 명령으로 등록 (LLM이 분석).

    # 2) style 경로 결정 (우선순위: --style > --template > 기본)
    style_path_override = None
    if args.style:
        style_path_override = args.style
        if not Path(style_path_override).exists():
            _user_error(
                f"style.yaml 없음: {style_path_override}",
                hint="--style 경로를 확인하거나 --template로 자동 분석을 활용하세요."
            )
    elif args.template:
        t = args.template
        if t.endswith(".pptx"):
            t = t[:-len(".pptx")]
        if "/" in t:
            base = t
        else:
            base = f"templates/{t}"
        pptx_candidate = f"{base}.pptx"
        if not Path(pptx_candidate).exists():
            _user_error(
                f"템플릿 .pptx 없음: {pptx_candidate}",
                hint=f"templates/ 폴더에 {Path(pptx_candidate).name} 파일을 넣고 다시 실행하세요."
            )
        style_path_override = f"{base}.style.yaml"
        # sync로 자동 생성될 거지만, 그래도 안 만들어졌으면 경고
        if not Path(style_path_override).exists():
            _user_warn(
                f"style.yaml 캐시 없음: {style_path_override}. "
                f"/onboard-template 명령으로 등록하세요 (LLM이 raw.yaml + 썸네일 보고 style.yaml 작성)."
            )

    # 3) style 로드
    try:
        if style_path_override:
            load_style(style_path_override)
        else:
            load_style()
    except yaml.YAMLError as e:
        _user_error(
            f"style.yaml 파싱 실패: {e}",
            hint="yaml 들여쓰기·따옴표·콜론을 확인하세요."
        )

    # 4) 콘텐츠 yaml 로드·검증
    try:
        data = yaml.safe_load(Path(src_yaml).read_text(encoding="utf-8"))
    except yaml.YAMLError as e:
        _user_error(
            f"콘텐츠 yaml 파싱 실패 ({src_yaml}): {e}",
            hint="yaml 들여쓰기·따옴표·콜론을 확인하세요. 한글 문자열은 작은따옴표로 감싸세요."
        )
    validate_content_yaml(data, src_yaml)

    if not Path(TEMPLATE_PATH).exists():
        print(f"warning: template not found at {TEMPLATE_PATH}", file=sys.stderr)
        return

    src_prs = Presentation(TEMPLATE_PATH)
    src_slides = list(src_prs.slides)

    out_prs = Presentation(TEMPLATE_PATH)
    xml_slides = out_prs.slides._sldIdLst
    for sld in list(xml_slides):
        xml_slides.remove(sld)

    print(f"template loaded: {TEMPLATE_PATH}")
    print(f"견본 슬라이드 개수: {len(src_slides)}")

    # 빌더는 yaml의 *장번호*를 직접 사용 (정책 없음). inherit은 ppt-designer agent의 책임.
    for sd in data.get("슬라이드", []):
        kind = sd.get("종류", "본문")
        title = sd.get("제목", "")
        # 종류 자동 재분류 (system_defaults.yaml kind_remapping)
        kind_eff = kind
        for rule in KIND_REMAPPING:
            if rule.get("from_kind") != kind:
                continue
            keywords = rule.get("title_contains_any", []) or []
            if any(kw in title for kw in keywords):
                kind_eff = rule.get("to_kind", kind)
                break

        # int|list|dict|chapter-keyed dict 모두 허용
        # chapter-keyed (2-level) dict면 slide.장번호로 1-level dict 직접 lookup
        raw_value = SAMPLE_MAP_1BASED.get(kind_eff)
        resolved = _resolve_chapter_variants(sd, raw_value)
        variants = _normalize_variants(resolved)
        sample_1based = pick_variant(sd, variants) if variants else None

        # 코드 fallback: layouts에서 null/None/0이면 견본 복제 안 하고 직접 그림
        use_code_fallback = (sample_1based is None) or (sample_1based == 0)

        if use_code_fallback:
            new_slide = make_blank_slide(out_prs)
            if kind_eff == "목차":
                draw_toc_slide(new_slide, sd, out_prs)
            elif kind_eff == "감사":
                # 간단한 감사 슬라이드 (제목만)
                slide_w = out_prs.slide_width / 914400
                slide_h = out_prs.slide_height / 914400
                box = new_slide.shapes.add_textbox(
                    Inches(0), Inches((slide_h - 1) / 2),
                    Inches(slide_w), Inches(1)
                )
                tf = box.text_frame
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                msg = sd.get("내용", {}).get("메시지", sd.get("제목", DEFAULTS["thanks_message"]))
                r = p.add_run()
                r.text = msg
                set_run_font(r, size=40, bold=True, color=NAVY)
            else:
                # 그 외 종류는 빈 슬라이드 + 제목만
                slide_w = out_prs.slide_width / 914400
                box = new_slide.shapes.add_textbox(
                    Inches(0.5), Inches(0.5),
                    Inches(slide_w - 1), Inches(0.8)
                )
                r = box.text_frame.paragraphs[0].add_run()
                r.text = sd.get("제목", "")
                set_run_font(r, size=24, bold=True, color=NAVY)
            add_notes(new_slide, sd.get("발표자_노트", ""))
            continue

        # 견본 복제 모드 (기존 동작)
        try:
            sample_idx_0 = int(sample_1based) - 1
        except (TypeError, ValueError):
            sample_idx_0 = 12  # "?" 등 미정 값 → 본문 fallback
        if sample_idx_0 < 0 or sample_idx_0 >= len(src_slides):
            sample_idx_0 = 12
        src_slide = src_slides[sample_idx_0]

        new_slide = duplicate_slide(out_prs, src_slide)

        # 그라데이션 mask → solid 흰 변환 (항상 적용, flag 무관 — LibreOffice 렌더 일관성)
        _solidify_gradient_masks(new_slide)
        # 디자인 가이드 픽처 제거 (항상 적용)
        _remove_decorative_pictures(new_slide)

        # *모든 placeholder text 일괄 비움* (font·색·shape 보존).
        # 견본의 placeholder 콘텐츠(견본 안내문 또는 발표자 콘텐츠) 무관하게 비움.
        # fill 함수가 매칭된 placeholder에 새 콘텐츠 채움. 매칭 안 된 placeholder는 빈 채로.
        for ph in list(new_slide.placeholders):
            if ph.has_text_frame and ph.text_frame.text.strip():
                set_run_text(ph, "", keep_font=True)

        # 견본 회사 콘텐츠 자동 제거 (옵션 활성화 시) — 비-placeholder만
        if SAMPLE_CLONE.get("clear_non_placeholder_text"):
            clear_non_placeholder_text(new_slide)

        # style.yaml의 sample_clone.remove_shape_coords 따라 *명시 좌표* shape 제거.
        # LLM (template-analyzer)이 vision으로 식별한 잔재 좌표. PY는 단순 lookup.
        _remove_listed_shapes(new_slide)

        if kind_eff == "표지":
            fill_cover(new_slide, sd)
        elif kind_eff == "목차":
            fill_toc(new_slide, sd, out_prs)
        elif kind_eff == "간지":
            fill_section_divider(new_slide, sd)
        elif kind_eff == "회사소개":
            # style.yaml의 company_intro.fill_function 으로 어느 fill 함수 쓸지 명시.
            # 정책은 yaml에, PY는 단순 dispatch.
            fill_fn_name = COMPANY_INTRO_FILL_FN or "fill_company_intro"
            if fill_fn_name == "fill_content":
                fill_content(new_slide, sd)
            else:
                fill_company_intro(new_slide, sd)
        elif kind_eff == "감사":
            fill_thanks(new_slide, sd)
        else:
            fill_content(new_slide, sd)

        # 챕터 라벨 (layout에 박힌 chap_no·chap_name 자리에 *콘텐츠 yaml의 장번호·장이름* 덮어쓰기)
        # 라벨별 enabled_for_kinds (yaml) 에 따라 그릴지 결정.
        draw_chapter_labels(new_slide, sd, kind=kind_eff)

        add_notes(new_slide, sd.get("발표자_노트", ""))

    out_prs.save(dst_pptx)
    print(f"saved: {dst_pptx}")

    # zip 중복 파일 제거 (LibreOffice 호환성)
    n_removed = dedupe_pptx(dst_pptx)
    if n_removed:
        print(f"dedupe: {n_removed} duplicate file(s) removed")


def dedupe_pptx(path):
    """pptx zip 내부 중복 파일 제거.
    같은 이름의 항목이 여러 개 있으면 마지막 위치의 항목만 유지.
    LibreOffice는 중복을 거부하므로 PowerPoint 호환 + LibreOffice 호환 동시 보장."""
    tmp = path + ".dedup.tmp"
    n_removed = 0
    with zipfile.ZipFile(path, 'r') as zin:
        infos = zin.infolist()
        last_idx = {}
        for i, info in enumerate(infos):
            last_idx[info.filename] = i

        with zipfile.ZipFile(tmp, 'w', zipfile.ZIP_DEFLATED) as zout:
            for i, info in enumerate(infos):
                if last_idx[info.filename] != i:
                    n_removed += 1
                    continue
                with zin.open(info) as src:
                    data = src.read()
                # ZipInfo를 재사용하면 압축 메타가 보존됨
                new_info = zipfile.ZipInfo(filename=info.filename,
                                            date_time=info.date_time)
                new_info.compress_type = zipfile.ZIP_DEFLATED
                new_info.external_attr = info.external_attr
                zout.writestr(new_info, data)

    os.replace(tmp, path)
    return n_removed


if __name__ == "__main__":
    main()
