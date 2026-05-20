#!/usr/bin/env python3
"""PPT 템플릿 *원시 데이터 추출* 도구.

용법:
    python3 scripts/analyze_template.py <template.pptx> --raw [output_raw.yaml]

PY는 원시 데이터 (layout 이름·shape 좌표·텍스트·테마)만 dump한다.
*의미 분석은 LLM (template-analyzer agent)이 100% 담당* — 챕터 구조·종류 매핑·
slot 좌표·dummy texts·색·폰트 모두 LLM이 raw.yaml + 썸네일 보고 직접 결정해
templates/<name>.style.yaml 작성.

이전의 점수 기반 자동 매핑 (legacy)은 회사·언어·디자인 시스템마다 일반화 안 되어
폐지됨. /onboard-template 명령 + template-analyzer agent 가 새 흐름.
"""
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent))
from pptx import Presentation
import yaml as pyyaml


# ─────────────────────────────────────────
# 유틸
# ─────────────────────────────────────────

def emu_to_inch(emu):
    if emu is None:
        return None
    return round(emu / 914400, 3)


# ─────────────────────────────────────────
# 테마 색·폰트 추출 (원시 — LLM이 의미 매핑)
# ─────────────────────────────────────────

def extract_theme_colors(prs):
    colors = {}
    try:
        master = prs.slide_masters[0]
        theme_part = None
        for rel in master.part.rels.values():
            if "theme" in rel.target_ref:
                theme_part = rel.target_part
                break
        if theme_part is None:
            return colors
        from lxml import etree
        tree = etree.fromstring(theme_part.blob)
        ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
        scheme = tree.find(".//a:clrScheme", ns)
        if scheme is None:
            return colors
        for child in scheme:
            tag = etree.QName(child.tag).localname
            srgb = child.find(".//a:srgbClr", ns)
            if srgb is not None:
                hex_val = srgb.get("val")
                if hex_val:
                    colors[tag] = f"#{hex_val.upper()}"
            else:
                sys_clr = child.find(".//a:sysClr", ns)
                if sys_clr is not None:
                    last = sys_clr.get("lastClr")
                    if last:
                        colors[tag] = f"#{last.upper()}"
    except Exception as e:
        print(f"warn: 테마 색상 추출 실패: {e}", file=sys.stderr)
    return colors


def extract_theme_fonts(prs):
    fonts = {}
    try:
        master = prs.slide_masters[0]
        theme_part = None
        for rel in master.part.rels.values():
            if "theme" in rel.target_ref:
                theme_part = rel.target_part
                break
        if theme_part is None:
            return fonts
        from lxml import etree
        tree = etree.fromstring(theme_part.blob)
        ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
        font_scheme = tree.find(".//a:fontScheme", ns)
        if font_scheme is None:
            return fonts
        for tag, path in [("major", ".//a:majorFont/a:latin"),
                          ("minor", ".//a:minorFont/a:latin"),
                          ("major_ea", ".//a:majorFont/a:ea"),
                          ("minor_ea", ".//a:minorFont/a:ea")]:
            el = font_scheme.find(path, ns)
            if el is not None and el.get("typeface"):
                fonts[tag] = el.get("typeface")
    except Exception as e:
        print(f"warn: 테마 폰트 추출 실패: {e}", file=sys.stderr)
    return fonts


# ─────────────────────────────────────────
# 슬라이드·shape 원시 정보 수집
# ─────────────────────────────────────────

def shape_info(shape):
    try:
        st = shape.shape_type
        shape_type_name = str(st).rsplit(".", 1)[-1] if st else None
    except (ValueError, NotImplementedError, AttributeError):
        shape_type_name = None
    info = {
        "name": shape.name,
        "is_placeholder": shape.is_placeholder,
        "shape_type": shape_type_name,
        "left": emu_to_inch(shape.left),
        "top": emu_to_inch(shape.top),
        "width": emu_to_inch(shape.width),
        "height": emu_to_inch(shape.height),
    }
    if shape.is_placeholder:
        info["ph_idx"] = shape.placeholder_format.idx
        info["ph_type"] = str(shape.placeholder_format.type)
    if shape.has_text_frame:
        info["text"] = shape.text_frame.text.strip()
        try:
            if shape.text_frame.paragraphs and shape.text_frame.paragraphs[0].runs:
                size = shape.text_frame.paragraphs[0].runs[0].font.size
                if size is not None:
                    info["font_pt"] = size.pt
        except Exception:
            pass
    return info


def analyze_layouts(prs):
    master = prs.slide_masters[0]
    layouts_info = []
    for idx, layout in enumerate(master.slide_layouts):
        placeholders = []
        for ph in layout.placeholders:
            placeholders.append({
                "idx": ph.placeholder_format.idx,
                "type": str(ph.placeholder_format.type),
                "name": ph.name,
                "left": emu_to_inch(ph.left),
                "top": emu_to_inch(ph.top),
                "width": emu_to_inch(ph.width),
                "height": emu_to_inch(ph.height),
            })
        extra_shapes = []
        for s in layout.shapes:
            if s.is_placeholder:
                continue
            extra_shapes.append(shape_info(s))
        layouts_info.append({
            "idx": idx,
            "name": layout.name,
            "placeholders": placeholders,
            "extras": extra_shapes,
        })
    return layouts_info


def _find_layout_idx(prs, layout):
    master = prs.slide_masters[0]
    for i, L in enumerate(master.slide_layouts):
        if L == layout:
            return i
    return -1


def analyze_samples(prs):
    samples = []
    for idx, slide in enumerate(prs.slides, start=1):
        shapes_info_list = [shape_info(s) for s in slide.shapes]
        samples.append({
            "idx": idx,
            "layout": slide.slide_layout.name,
            "layout_idx": _find_layout_idx(prs, slide.slide_layout),
            "shapes": shapes_info_list,
        })
    return samples


# ─────────────────────────────────────────
# 원시 데이터 dump (LLM 분석용)
# ─────────────────────────────────────────

def dump_raw(pptx_path, out_path):
    """layout·slide의 모든 shape 정보를 yaml로 dump. *의미 분석 없음*.
    LLM (template-analyzer agent)이 이 yaml + 썸네일 보고 의미 매핑한 style.yaml 작성.
    """
    prs = Presentation(pptx_path)
    print(f"원시 추출 중: {pptx_path}", file=sys.stderr)

    slide_w = emu_to_inch(prs.slide_width)
    slide_h = emu_to_inch(prs.slide_height)
    layouts = analyze_layouts(prs)
    samples = analyze_samples(prs)
    theme_colors = extract_theme_colors(prs)
    theme_fonts = extract_theme_fonts(prs)

    output = {
        "_generated_by": "scripts/analyze_template.py --raw",
        "_template_source": pptx_path,
        "_purpose": "원시 데이터. LLM (template-analyzer agent)이 이 yaml과 thumbnails 보고 의미 분석.",
        "slide_size": {
            "width_in": round(slide_w or 13.33, 3),
            "height_in": round(slide_h or 7.5, 3),
        },
        "theme": {
            "colors": theme_colors,
            "fonts": theme_fonts,
        },
        "layouts": [
            {
                "idx": L["idx"],
                "name": L["name"],
                "placeholders": L["placeholders"],
                "extras": L["extras"],
            }
            for L in layouts
        ],
        "slides": [
            {
                "idx": S["idx"],
                "layout_name": S["layout"],
                "layout_idx": S["layout_idx"],
                "shapes": S["shapes"],
            }
            for S in samples
        ],
    }

    with open(out_path, "w", encoding="utf-8") as f:
        pyyaml.safe_dump(output, f, allow_unicode=True, sort_keys=False, default_flow_style=False)
    print(f"저장: {out_path}", file=sys.stderr)
    print(f"  layouts: {len(layouts)}개", file=sys.stderr)
    print(f"  slides: {len(samples)}개", file=sys.stderr)
    print(f"  → LLM이 이 raw.yaml + thumbnails 보고 style.yaml 작성", file=sys.stderr)


def main():
    if len(sys.argv) < 2 or "--help" in sys.argv or "-h" in sys.argv:
        print("사용:")
        print("  python3 scripts/analyze_template.py <template.pptx> --raw [output_raw.yaml]")
        print()
        print("원시 데이터 (layouts, slides, shapes, theme) 만 추출. 의미 매핑은 LLM이 담당.")
        print("자세한 흐름: /onboard-template 슬래시 명령 + .claude/agents/template-analyzer.md")
        sys.exit(1)

    pptx_path = sys.argv[1]

    if len(sys.argv) >= 3 and sys.argv[2] == "--raw":
        raw_out = sys.argv[3] if len(sys.argv) >= 4 else pptx_path[:-len(".pptx")] + ".raw.yaml"
        return dump_raw(pptx_path, raw_out)

    # legacy 모드 호출 시 명확한 에러 + 안내
    print("ERROR: legacy 자동 매핑 모드는 폐지됨.", file=sys.stderr)
    print("  PY가 의미 분석을 하지 않습니다. 대신 LLM (template-analyzer agent)이 담당.", file=sys.stderr)
    print(file=sys.stderr)
    print("  사용:", file=sys.stderr)
    print(f"    python3 scripts/analyze_template.py {pptx_path} --raw", file=sys.stderr)
    print(f"    # 또는 Claude Code에서: /onboard-template", file=sys.stderr)
    sys.exit(2)


if __name__ == "__main__":
    main()
