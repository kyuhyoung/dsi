"""슬라이드 구성안 YAML → .pptx 변환기

DSI 시스템 smoke test용 간이 변환기.
실제 운영에서는 공식 pptx skill 활용 권장.
"""

import sys
import yaml
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


def add_title_slide(prs, info):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = info.get('사업명', '제안서')
    sub = slide.placeholders[1]
    lines = []
    if '발주처' in info:
        lines.append(f"발주처: {info['발주처']}")
    if '제안사' in info:
        lines.append(f"제안사: {info['제안사']}")
    if '제출일' in info:
        lines.append(f"제출일: {info['제출일']}")
    sub.text = '\n'.join(lines)


def add_content_slide(prs, slide_def):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = f"{slide_def.get('번호', '')}  {slide_def.get('제목', '')}"

    tf = slide.placeholders[1].text_frame
    tf.clear()

    # 핵심 메시지
    msgs = slide_def.get('핵심메시지', [])
    if msgs:
        for i, m in enumerate(msgs):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"• {m}"
            p.font.size = Pt(18)

    # 시각 요소 안내
    visuals = slide_def.get('시각요소', [])
    if visuals:
        p = tf.add_paragraph()
        p.text = ""
        for v in visuals:
            p2 = tf.add_paragraph()
            p2.text = f"[시각] {v}"
            p2.font.size = Pt(12)
            p2.font.color.rgb = RGBColor(0x80, 0x80, 0x80)
            p2.font.italic = True

    # 발표자 노트
    note = slide_def.get('발표자_노트', '')
    if note:
        slide.notes_slide.notes_text_frame.text = note

    # 내용 dict (표지·목차·결론·QA 등)
    content = slide_def.get('내용', {})
    if content and not msgs:
        for i, (k, v) in enumerate(content.items()):
            if isinstance(v, list):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = f"{k}:"
                p.font.size = Pt(18)
                p.font.bold = True
                for item in v:
                    p_item = tf.add_paragraph()
                    p_item.text = f"  • {item}"
                    p_item.font.size = Pt(14)
            else:
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = f"{k}: {v}"
                p.font.size = Pt(16)


def yaml_to_pptx(yaml_path, pptx_path):
    with open(yaml_path, 'r', encoding='utf-8') as f:
        data = yaml.safe_load(f)

    prs = Presentation()
    pres_info = data.get('프레젠테이션', {})
    slides_def = data.get('슬라이드', [])

    for slide_def in slides_def:
        종류 = slide_def.get('종류', '본문')
        if 종류 == '표지':
            add_title_slide(prs, slide_def.get('내용', {}))
        else:
            add_content_slide(prs, slide_def)

    prs.save(pptx_path)
    print(f"✓ 생성: {pptx_path}")
    print(f"  슬라이드 수: {len(slides_def)}")


if __name__ == '__main__':
    yaml_to_pptx(sys.argv[1], sys.argv[2])
