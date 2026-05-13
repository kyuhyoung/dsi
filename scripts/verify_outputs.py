"""DSI smoke test 산출물 무결성 검증 스크립트

사용법:
    python3 scripts/verify_outputs.py output/20260513
"""

import sys
from pathlib import Path

try:
    from docx import Document
    from pptx import Presentation
    import yaml
except ImportError as e:
    print(f"의존성 누락: {e}. './install.sh' 먼저 실행.")
    sys.exit(2)


def verify(workspace: Path) -> int:
    errors = []

    # 1. analysis.yaml
    print("=== analysis.yaml ===")
    analysis = next(workspace.glob('analysis*.yaml'), None)
    if not analysis:
        print("  ✗ analysis.yaml 없음")
        errors.append('analysis.yaml missing')
    else:
        try:
            with open(analysis) as f:
                data = yaml.safe_load(f)
            assert '사업개요' in data, '사업개요 키 누락'
            assert '평가배점' in data, '평가배점 키 누락'
            print(f"  ✓ 파싱 정상, 평가 총점 {data['평가배점'].get('총점', '?')}점")
            print(f"  ✓ 누락 항목 {len(data.get('누락_항목', []))} 개")
            print(f"  ✓ 위험 신호 {len(data.get('위험_신호', []))} 개")
        except Exception as e:
            errors.append(f"analysis.yaml: {e}")
            print(f"  ✗ {e}")

    # 2. proposal .md
    print("\n=== proposal .md ===")
    md_files = list(workspace.glob('proposal_*.md'))
    if not md_files:
        print("  ✗ proposal*.md 없음")
        errors.append('proposal md missing')
    else:
        md = md_files[0].read_text(encoding='utf-8')
        sections = md.count('## ')
        print(f"  ✓ {len(md)} 글자, 섹션 {sections} 개")
        if '당사는' in md:
            print("  ✓ 격식체 사용 ('당사는')")
        else:
            print("  ⚠ 격식체 미검출")

    # 3. proposal .docx
    print("\n=== proposal .docx ===")
    docx_files = list(workspace.glob('proposal_*.docx'))
    if not docx_files:
        print("  ✗ proposal*.docx 없음")
        errors.append('proposal docx missing')
    else:
        try:
            doc = Document(docx_files[0])
            korean = sum(1 for p in doc.paragraphs
                         if any('가' <= c <= '힯' for c in p.text))
            print(f"  ✓ 단락 {len(doc.paragraphs)} 개, 표 {len(doc.tables)} 개")
            print(f"  ✓ 한국어 단락 {korean} 개")
        except Exception as e:
            errors.append(f"proposal.docx: {e}")
            print(f"  ✗ {e}")

    # 4. slides .yaml
    print("\n=== slides .yaml ===")
    sy = next(workspace.glob('slides_*.yaml'), None)
    if not sy:
        print("  ✗ slides*.yaml 없음")
        errors.append('slides yaml missing')
    else:
        try:
            with open(sy) as f:
                data = yaml.safe_load(f)
            slides = data.get('슬라이드', [])
            notes = sum(1 for s in slides if s.get('발표자_노트'))
            print(f"  ✓ 슬라이드 {len(slides)} 개, 발표자 노트 {notes} 개")
        except Exception as e:
            errors.append(f"slides.yaml: {e}")
            print(f"  ✗ {e}")

    # 5. .pptx
    print("\n=== presentation .pptx ===")
    pptx_files = list(workspace.glob('*.pptx'))
    if not pptx_files:
        print("  ✗ .pptx 없음")
        errors.append('pptx missing')
    else:
        try:
            prs = Presentation(pptx_files[0])
            slides = list(prs.slides)
            notes_count = sum(1 for s in slides
                              if s.has_notes_slide
                              and s.notes_slide.notes_text_frame.text)
            print(f"  ✓ 슬라이드 {len(slides)} 개")
            print(f"  ✓ 발표자 노트 포함 {notes_count} 개")
        except Exception as e:
            errors.append(f"pptx: {e}")
            print(f"  ✗ {e}")

    print("\n=== 종합 결과 ===")
    if errors:
        print(f"  ✗ {len(errors)} 개 오류:")
        for e in errors:
            print(f"    - {e}")
        return 1
    print("  ✓ 모든 산출물 검증 통과")
    return 0


if __name__ == '__main__':
    if len(sys.argv) != 2:
        print(f"사용법: python3 {sys.argv[0]} <output_dir>")
        sys.exit(2)
    ws = Path(sys.argv[1])
    if not ws.is_dir():
        print(f"디렉토리 없음: {ws}")
        sys.exit(2)
    sys.exit(verify(ws))
